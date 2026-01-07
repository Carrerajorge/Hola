#!/usr/bin/env python3
"""
╔══════════════════════════════════════════════════════════════════════════════╗
║            🚀 SISTEMA COMPLETO DE AGENTE IA v2.0 🚀                         ║
╠══════════════════════════════════════════════════════════════════════════════╣
║  ✨ PowerPoint, Word, Excel  |  🔍 Búsqueda Web  |  🌐 Navegación           ║
║  📁 Gestión de Archivos      |  🐍 Python       |  🔬 Investigación         ║
║  👁️ Visualización en Tiempo Real del Proceso Completo                       ║
╚══════════════════════════════════════════════════════════════════════════════╝
"""

import os, sys, re, json, shutil, hashlib, asyncio, tempfile, subprocess, shlex
import mimetypes, time, logging, uuid, traceback
from pathlib import Path
from datetime import datetime
from typing import Dict, Any, Optional, List, Union
from dataclasses import dataclass, field
from enum import Enum
from abc import ABC, abstractmethod
from collections import deque, defaultdict
import inspect

# ═══════════════════════════════════════════════════════════════════════════════
# INSTALACIÓN DE DEPENDENCIAS
# ═══════════════════════════════════════════════════════════════════════════════

def _install(pkg):
    # Validate package name to ensure it only contains safe characters
    if not re.match(r'^[a-zA-Z0-9._-]+$', pkg):
        raise ValueError(f"Invalid package name: {pkg}")
    subprocess.run([sys.executable, "-m", "pip", "install", pkg, "-q"], capture_output=True)

PKGS = [("aiofiles","aiofiles"),("rich","rich"),("httpx","httpx"),("beautifulsoup4","bs4"),
        ("python-pptx","pptx"),("python-docx","docx"),("openpyxl","openpyxl")]

print("🔧 Verificando dependencias...")
for pip_n, imp_n in PKGS:
    try: __import__(imp_n); print(f"  ✅ {pip_n}")
    except:
        print(f"  📦 {pip_n}...", end=" ", flush=True)
        _install(pip_n)
        try: __import__(imp_n); print("✅")
        except: print("❌")
print()

import aiofiles
from rich.console import Console
from rich.panel import Panel
from rich.table import Table
from rich.tree import Tree
from rich.rule import Rule
from rich.prompt import Prompt
from rich import box

console = Console()
logging.basicConfig(level=logging.WARNING)

# Opcionales
try: import httpx; HTTPX_OK = True
except: HTTPX_OK = False
try: from bs4 import BeautifulSoup; BS4_OK = True
except: BS4_OK = False
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    PPTX_OK = True
except: PPTX_OK = False
try:
    from docx import Document as DocxDoc
    from docx.shared import Inches as DInches, Pt as DPt
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    DOCX_OK = True
except: DOCX_OK = False
try:
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    from openpyxl.chart import BarChart, Reference
    XLSX_OK = True
except: XLSX_OK = False

# ═══════════════════════════════════════════════════════════════════════════════
# ENUMERACIONES Y DATACLASSES
# ═══════════════════════════════════════════════════════════════════════════════

class ThreatLevel(Enum):
    SAFE="safe"; LOW="low"; MEDIUM="medium"; HIGH="high"; CRITICAL="critical"

class SecurityAction(Enum):
    ALLOW="allow"; WARN="warn"; BLOCK="block"; LOG_AND_BLOCK="log_and_block"

class ExecutionStatus(Enum):
    PENDING="pending"; RUNNING="running"; COMPLETED="completed"
    FAILED="failed"; TIMEOUT="timeout"; BLOCKED="blocked"

class AgentState(Enum):
    IDLE="idle"; ANALYZING="analyzing"; PLANNING="planning"
    EXECUTING="executing"; DELIVERING="delivering"; ERROR="error"

class PhaseStatus(Enum):
    PENDING="pending"; IN_PROGRESS="in_progress"; COMPLETED="completed"; FAILED="failed"

class ToolCategory(Enum):
    SYSTEM="system"; FILE="file"; DOCUMENT="document"; SEARCH="search"
    BROWSER="browser"; COMMUNICATION="communication"; DEVELOPMENT="development"

@dataclass
class SecurityAnalysis:
    command: str; is_safe: bool; threat_level: ThreatLevel
    action: SecurityAction; warnings: List[str] = field(default_factory=list)

@dataclass
class ExecutionResult:
    command: str; status: ExecutionStatus; return_code: int = None
    stdout: str = ""; stderr: str = ""; execution_time: float = 0.0
    error_message: str = ""
    @property
    def success(self): return self.status == ExecutionStatus.COMPLETED and self.return_code == 0

@dataclass
class FileOperationResult:
    success: bool; operation: str; path: str
    message: str = ""; data: Any = None; error: str = None

@dataclass
class ToolResult:
    success: bool; tool_name: str; data: Any = None
    message: str = ""; error: str = None; execution_time: float = 0.0
    files_created: List[str] = field(default_factory=list)

@dataclass
class Step:
    id: str; description: str; tool: str; params: Dict[str, Any]
    status: PhaseStatus = PhaseStatus.PENDING; result: Any = None; error: str = None
    def complete(self, r): self.status = PhaseStatus.COMPLETED; self.result = r
    def fail(self, e): self.status = PhaseStatus.FAILED; self.error = e

@dataclass
class Phase:
    id: str; name: str; description: str; icon: str = "📋"
    steps: List[Step] = field(default_factory=list)
    status: PhaseStatus = PhaseStatus.PENDING
    @property
    def pending_steps(self): return [s for s in self.steps if s.status == PhaseStatus.PENDING]
    @property
    def is_complete(self): return all(s.status == PhaseStatus.COMPLETED for s in self.steps)
    def get_next_step(self): return self.pending_steps[0] if self.pending_steps else None

@dataclass
class TaskPlan:
    task_id: str; objective: str
    phases: List[Phase] = field(default_factory=list)
    current_phase_index: int = 0
    @property
    def is_complete(self): return all(p.is_complete for p in self.phases)
    @property
    def progress(self):
        total = sum(len(p.steps) for p in self.phases)
        done = sum(1 for p in self.phases for s in p.steps if s.status == PhaseStatus.COMPLETED)
        return (done/total)*100 if total > 0 else 0
    def get_current_phase(self):
        return self.phases[self.current_phase_index] if 0 <= self.current_phase_index < len(self.phases) else None
    def advance(self):
        cur = self.get_current_phase()
        if cur and cur.is_complete:
            cur.status = PhaseStatus.COMPLETED
            self.current_phase_index += 1
            nxt = self.get_current_phase()
            if nxt: nxt.status = PhaseStatus.IN_PROGRESS

DEFAULT_WORKSPACE = os.path.join(os.getcwd(), "workspace_v2")
DEFAULT_OUTPUT = os.path.join(DEFAULT_WORKSPACE, "output")

# ═══════════════════════════════════════════════════════════════════════════════
# VISUALIZADOR EN TIEMPO REAL
# ═══════════════════════════════════════════════════════════════════════════════

class Visualizer:
    def __init__(self):
        self.stats = {'tools': 0, 'files': 0, 'searches': 0, 'errors': 0, 'start': None}
    
    def start(self): self.stats['start'] = datetime.now()
    def log(self, msg, icon="ℹ️"): console.print(f"    {icon} {msg}")
    def phase(self, name, icon="📋"): console.print(f"\n  [bold blue]{icon} {name}[/bold blue]")
    def step(self, desc): console.print(f"    [cyan]→[/cyan] {desc}")
    def tool(self, name, params):
        p = ", ".join(f"{k}={repr(v)[:25]}" for k,v in list(params.items())[:2])
        console.print(f"      [yellow]🔧 {name}[/yellow]({p})")
        self.stats['tools'] += 1
    def result(self, ok, msg):
        if ok: console.print(f"        [green]✓[/green] {msg[:80]}")
        else: console.print(f"        [red]✗[/red] {msg[:80]}"); self.stats['errors'] += 1
    def file_created(self, path):
        console.print(f"      [green]📄 Creado:[/green] {path}")
        self.stats['files'] += 1
    def research(self, q, n):
        console.print(f"      [cyan]🔬[/cyan] {q} → {n} resultados")
        self.stats['searches'] += 1
    def summary(self, plan):
        elapsed = (datetime.now() - self.stats['start']).total_seconds() if self.stats['start'] else 0
        console.print(Rule("[bold green]✅ COMPLETADO[/bold green]"))
        t = Table(show_header=False, box=box.SIMPLE)
        t.add_column("", style="cyan"); t.add_column("", style="green")
        t.add_row("⏱️ Duración", f"{elapsed:.1f}s")
        t.add_row("📊 Progreso", f"{plan.progress:.0f}%")
        t.add_row("🔧 Herramientas", str(self.stats['tools']))
        t.add_row("📄 Archivos", str(self.stats['files']))
        t.add_row("🔍 Búsquedas", str(self.stats['searches']))
        console.print(t)

viz = Visualizer()

# ═══════════════════════════════════════════════════════════════════════════════
# SECURITY GUARD
# ═══════════════════════════════════════════════════════════════════════════════

class SecurityGuard:
    CRITICAL = [r"rm\s+(-[rfv]+\s+)*/?$", r"rm\s+(-[rfv]+\s+)*/\*", r":\(\)\s*\{\s*:\|:\s*&",
                r"mkfs\.", r"dd\s+if=/dev", r"curl.*\|\s*(bash|sh)", r"chmod\s+777\s+/"]
    SAFE_CMDS = {'ls','pwd','cd','cat','echo','date','whoami','python3','pip','node','npm',
                 'git','curl','mkdir','touch','cp','mv','head','tail','grep','find','wc'}
    
    def __init__(self, root=None):
        self.sandbox_root = Path(root or DEFAULT_WORKSPACE).resolve()
        self.sandbox_root.mkdir(parents=True, exist_ok=True)
        self._critical = [re.compile(p, re.I) for p in self.CRITICAL]
    
    def analyze(self, cmd):
        cmd = cmd.strip()
        if not cmd: return SecurityAnalysis(cmd, True, ThreatLevel.SAFE, SecurityAction.ALLOW)
        for p in self._critical:
            if p.search(cmd):
                return SecurityAnalysis(cmd, False, ThreatLevel.CRITICAL, SecurityAction.LOG_AND_BLOCK, ["⚠️ BLOQUEADO"])
        base = cmd.split()[0].split('/')[-1] if cmd.split() else ""
        if base in self.SAFE_CMDS:
            return SecurityAnalysis(cmd, True, ThreatLevel.SAFE, SecurityAction.ALLOW)
        return SecurityAnalysis(cmd, True, ThreatLevel.LOW, SecurityAction.ALLOW)
    
    def validate_path(self, path):
        try:
            resolved = (self.sandbox_root / path).resolve() if not os.path.isabs(path) else Path(path).resolve()
            try:
                resolved.relative_to(self.sandbox_root)
                return {'allowed': True, 'path': str(resolved)}
            except: return {'allowed': False, 'reason': 'Fuera del sandbox'}
        except Exception as e: return {'allowed': False, 'reason': str(e)}

# ═══════════════════════════════════════════════════════════════════════════════
# COMMAND EXECUTOR Y FILE MANAGER
# ═══════════════════════════════════════════════════════════════════════════════

class CommandExecutor:
    def __init__(self, security=None, timeout=30):
        self.security = security or SecurityGuard()
        self.timeout = timeout
        self.workdir = self.security.sandbox_root
    
    async def execute(self, cmd, timeout=None):
        start = time.time()
        analysis = self.security.analyze(cmd)
        if analysis.action == SecurityAction.LOG_AND_BLOCK:
            return ExecutionResult(cmd, ExecutionStatus.BLOCKED, error_message="Bloqueado")
        try:
            # For simple commands without shell features, use safer exec mode
            # For complex commands with pipes/redirects, continue using shell but with awareness
            if not any(char in cmd for char in ['|', '>', '<', '&', ';', '$', '`', '(', ')']):
                # Simple command - parse safely
                try:
                    args = shlex.split(cmd)
                    proc = await asyncio.create_subprocess_exec(
                        *args, stdout=asyncio.subprocess.PIPE, stderr=asyncio.subprocess.PIPE,
                        cwd=str(self.workdir), env=os.environ.copy())
                except ValueError as e:
                    # Reject malformed commands instead of falling back to shell
                    return ExecutionResult(cmd, ExecutionStatus.FAILED, 
                        error_message=f"Comando malformado: {str(e)}")
            else:
                # Complex command with shell features - use shell but it's been security-analyzed
                proc = await asyncio.create_subprocess_shell(
                    cmd, stdout=asyncio.subprocess.PIPE, stderr=asyncio.subprocess.PIPE,
                    cwd=str(self.workdir), env=os.environ.copy())
            stdout, stderr = await asyncio.wait_for(proc.communicate(), timeout or self.timeout)
            return ExecutionResult(cmd, ExecutionStatus.COMPLETED, proc.returncode,
                stdout.decode('utf-8', errors='replace'), stderr.decode('utf-8', errors='replace'),
                time.time()-start)
        except asyncio.TimeoutError:
            return ExecutionResult(cmd, ExecutionStatus.TIMEOUT, error_message="Timeout")
        except Exception as e:
            return ExecutionResult(cmd, ExecutionStatus.FAILED, error_message=str(e))
    
    async def run_script(self, code, interpreter="python3", timeout=None):
        # Whitelist validation to prevent command injection
        allowed_interpreters = {'python', 'python3', 'python3.8', 'python3.9', 'python3.10', 'python3.11', 'python3.12'}
        if interpreter not in allowed_interpreters:
            return ExecutionResult(f"{interpreter} (rejected)", ExecutionStatus.FAILED, 
                error_message=f"Interpreter no permitido: {interpreter}")
        path = self.workdir / f"_script_{uuid.uuid4().hex[:6]}.py"
        try:
            async with aiofiles.open(path, 'w') as f: await f.write(code)
            os.chmod(path, 0o755)
            return await self.execute(f"{interpreter} {path}", timeout)
        finally:
            try: path.unlink()
            except: pass


class FileManager:
    def __init__(self, security=None):
        self.security = security or SecurityGuard()
        self.root = self.security.sandbox_root
    
    def _resolve(self, p): return (self.root / p).resolve() if not os.path.isabs(p) else Path(p).resolve()
    def _validate(self, p): return self.security.validate_path(str(self._resolve(p)))
    
    async def read(self, path):
        v = self._validate(path)
        if not v.get('allowed'): return FileOperationResult(False, 'read', path, error=v.get('reason'))
        try:
            async with aiofiles.open(self._resolve(path), 'r') as f:
                return FileOperationResult(True, 'read', path, data=await f.read())
        except FileNotFoundError: return FileOperationResult(False, 'read', path, error="No encontrado")
        except Exception as e: return FileOperationResult(False, 'read', path, error=str(e))
    
    async def write(self, path, content, create_dirs=True):
        v = self._validate(path)
        if not v.get('allowed'): return FileOperationResult(False, 'write', path, error=v.get('reason'))
        resolved = self._resolve(path)
        try:
            if create_dirs: resolved.parent.mkdir(parents=True, exist_ok=True)
            async with aiofiles.open(resolved, 'w') as f: await f.write(content)
            return FileOperationResult(True, 'write', str(resolved), message=f"{len(content)} bytes")
        except Exception as e: return FileOperationResult(False, 'write', path, error=str(e))
    
    async def delete(self, path):
        v = self._validate(path)
        if not v.get('allowed'): return FileOperationResult(False, 'delete', path, error=v.get('reason'))
        resolved = self._resolve(path)
        try:
            if resolved.is_file(): resolved.unlink()
            elif resolved.is_dir(): shutil.rmtree(resolved)
            return FileOperationResult(True, 'delete', str(resolved))
        except Exception as e: return FileOperationResult(False, 'delete', path, error=str(e))
    
    async def list_dir(self, path="."):
        v = self._validate(path)
        if not v.get('allowed'): return FileOperationResult(False, 'list', path, error=v.get('reason'))
        resolved = self._resolve(path)
        try:
            items = [{'name': p.name, 'is_dir': p.is_dir(), 'size': p.stat().st_size if p.is_file() else 0}
                     for p in resolved.glob('*') if not p.name.startswith('.')]
            return FileOperationResult(True, 'list', str(resolved), data={'items': items, 'count': len(items)})
        except Exception as e: return FileOperationResult(False, 'list', path, error=str(e))
    
    async def mkdir(self, path):
        v = self._validate(path)
        if not v.get('allowed'): return FileOperationResult(False, 'mkdir', path, error=v.get('reason'))
        try:
            self._resolve(path).mkdir(parents=True, exist_ok=True)
            return FileOperationResult(True, 'mkdir', path)
        except Exception as e: return FileOperationResult(False, 'mkdir', path, error=str(e))
    
    async def exists(self, path):
        v = self._validate(path)
        return v.get('allowed') and self._resolve(path).exists()


# ═══════════════════════════════════════════════════════════════════════════════
# DOCUMENT CREATOR - PPTX, DOCX, XLSX
# ═══════════════════════════════════════════════════════════════════════════════

class DocumentCreator:
    def __init__(self, output_dir=None):
        self.output_dir = Path(output_dir or DEFAULT_OUTPUT)
        self.output_dir.mkdir(parents=True, exist_ok=True)
    
    async def create_pptx(self, title, slides, theme="professional", filename=None):
        if not PPTX_OK: return ToolResult(False, "pptx", error="python-pptx no instalado")
        try:
            prs = Presentation()
            prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
            colors = {'professional': RGBColor(46,134,171), 'modern': RGBColor(52,152,219)}
            color = colors.get(theme, colors['professional'])
            
            # Slide título
            sl = prs.slides.add_slide(prs.slide_layouts[6])
            bg = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
            bg.fill.solid(); bg.fill.fore_color.rgb = color; bg.line.fill.background()
            tb = sl.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.3), Inches(1.5))
            tf = tb.text_frame; p = tf.paragraphs[0]
            p.text = title; p.font.size = Pt(54); p.font.bold = True
            p.font.color.rgb = RGBColor(255,255,255); p.alignment = PP_ALIGN.CENTER
            
            # Slides contenido
            for i, sd in enumerate(slides):
                sl = prs.slides.add_slide(prs.slide_layouts[6])
                bar = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.15))
                bar.fill.solid(); bar.fill.fore_color.rgb = color; bar.line.fill.background()
                
                tb = sl.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.3), Inches(0.8))
                tf = tb.text_frame; p = tf.paragraphs[0]
                p.text = sd.get('title', f'Slide {i+2}'); p.font.size = Pt(36); p.font.bold = True
                
                if sd.get('bullets'):
                    for j, b in enumerate(sd['bullets'][:8]):
                        bb = sl.shapes.add_textbox(Inches(0.7), Inches(1.4+j*0.6), Inches(11.5), Inches(0.5))
                        bf = bb.text_frame; bp = bf.paragraphs[0]
                        bp.text = f"• {b}"; bp.font.size = Pt(22)
                elif sd.get('content'):
                    cb = sl.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(12.3), Inches(5))
                    cf = cb.text_frame; cf.word_wrap = True
                    cp = cf.paragraphs[0]; cp.text = sd['content']; cp.font.size = Pt(20)
            
            fn = filename or f"presentacion_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
            fp = self.output_dir / fn
            prs.save(str(fp))
            return ToolResult(True, "pptx", data={'path': str(fp)}, message=f"Creado: {fn}", files_created=[str(fp)])
        except Exception as e: return ToolResult(False, "pptx", error=str(e))
    
    async def create_docx(self, title, content, author="Agente IA", filename=None):
        if not DOCX_OK: return ToolResult(False, "docx", error="python-docx no instalado")
        try:
            doc = DocxDoc()
            doc.add_heading(title, 0).alignment = WD_ALIGN_PARAGRAPH.CENTER
            doc.add_paragraph("─" * 50).alignment = WD_ALIGN_PARAGRAPH.CENTER
            m = doc.add_paragraph(); m.alignment = WD_ALIGN_PARAGRAPH.CENTER
            m.add_run(f"Autor: {author} | {datetime.now().strftime('%d/%m/%Y')}").italic = True
            doc.add_paragraph()
            
            if isinstance(content, str):
                for para in content.split('\n\n'):
                    if para.strip(): doc.add_paragraph(para.strip())
            else:
                for sec in content:
                    if sec.get('title'): doc.add_heading(sec['title'], level=sec.get('level', 1))
                    if sec.get('content'): doc.add_paragraph(sec['content'])
                    for b in sec.get('bullets', []): doc.add_paragraph(b, style='List Bullet')
            
            fn = filename or f"documento_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
            fp = self.output_dir / fn
            doc.save(str(fp))
            return ToolResult(True, "docx", data={'path': str(fp)}, message=f"Creado: {fn}", files_created=[str(fp)])
        except Exception as e: return ToolResult(False, "docx", error=str(e))
    
    async def create_xlsx(self, title, sheets, filename=None):
        if not XLSX_OK: return ToolResult(False, "xlsx", error="openpyxl no instalado")
        try:
            wb = Workbook(); wb.remove(wb.active)
            hf = Font(bold=True, color="FFFFFF"); hfill = PatternFill("solid", fgColor="2E86AB")
            border = Border(left=Side('thin'), right=Side('thin'), top=Side('thin'), bottom=Side('thin'))
            
            for sd in sheets:
                ws = wb.create_sheet(title=sd.get('name', 'Hoja')[:31])
                headers = sd.get('headers', []); rows = sd.get('rows', [])
                
                for c, h in enumerate(headers, 1):
                    cell = ws.cell(1, c, h)
                    cell.font = hf; cell.fill = hfill; cell.border = border
                    ws.column_dimensions[cell.column_letter].width = max(15, len(str(h))+5)
                
                for ri, row in enumerate(rows, 2):
                    for ci, val in enumerate(row, 1):
                        cell = ws.cell(ri, ci, val); cell.border = border
                
                if sd.get('chart') and headers and rows:
                    chart = BarChart(); chart.title = sd['chart'].get('title', 'Gráfico')
                    data = Reference(ws, min_col=2, min_row=1, max_col=len(headers), max_row=len(rows)+1)
                    cats = Reference(ws, min_col=1, min_row=2, max_row=len(rows)+1)
                    chart.add_data(data, titles_from_data=True); chart.set_categories(cats)
                    ws.add_chart(chart, f"A{len(rows)+4}")
                
                ws.freeze_panes = 'A2'
            
            fn = filename or f"excel_{datetime.now().strftime('%Y%m%d_%H%M%S')}.xlsx"
            fp = self.output_dir / fn
            wb.save(str(fp))
            return ToolResult(True, "xlsx", data={'path': str(fp)}, message=f"Creado: {fn}", files_created=[str(fp)])
        except Exception as e: return ToolResult(False, "xlsx", error=str(e))

# ═══════════════════════════════════════════════════════════════════════════════
# WEB SEARCH Y BROWSER
# ═══════════════════════════════════════════════════════════════════════════════

class WebSearchEngine:
    def __init__(self): self.cache = {}
    
    async def search(self, query, num=10):
        if query in self.cache: return self.cache[query]
        if not HTTPX_OK: return {'query': query, 'results': [], 'error': 'httpx no disponible'}
        try:
            async with httpx.AsyncClient(timeout=15) as client:
                r = await client.get("https://api.duckduckgo.com/", params={'q': query, 'format': 'json', 'no_html': '1'})
                data = r.json(); results = []
                if data.get('AbstractText'):
                    results.append({'title': data.get('Heading', query), 'snippet': data['AbstractText'][:300],
                                    'url': data.get('AbstractURL', '')})
                for t in data.get('RelatedTopics', [])[:num-1]:
                    if isinstance(t, dict) and t.get('Text'):
                        results.append({'title': t['Text'][:80], 'snippet': t['Text'], 'url': t.get('FirstURL', '')})
                res = {'query': query, 'results': results[:num], 'total': len(results)}
                self.cache[query] = res
                return res
        except Exception as e: return {'query': query, 'results': [], 'error': str(e)}


class WebBrowser:
    async def fetch(self, url, extract_text=True):
        if not HTTPX_OK: return {'url': url, 'error': 'httpx no disponible'}
        try:
            async with httpx.AsyncClient(timeout=15, follow_redirects=True) as client:
                r = await client.get(url)
                content = r.text; title = ""
                if extract_text and BS4_OK:
                    soup = BeautifulSoup(content, 'html.parser')
                    title = soup.title.text.strip() if soup.title else ""
                    for tag in soup(['script', 'style', 'nav', 'footer']): tag.decompose()
                    content = soup.get_text('\n', strip=True)[:8000]
                return {'url': str(r.url), 'title': title, 'content': content, 'status': r.status_code}
        except Exception as e: return {'url': url, 'error': str(e)}


# ═══════════════════════════════════════════════════════════════════════════════
# HERRAMIENTAS (TOOLS)
# ═══════════════════════════════════════════════════════════════════════════════

class BaseTool(ABC):
    name = "base"; description = "Base"; category = ToolCategory.SYSTEM
    def __init__(self): self.enabled = True; self.count = 0
    @abstractmethod
    async def execute(self, **kw) -> ToolResult: pass
    async def __call__(self, **kw):
        start = time.time()
        try:
            viz.tool(self.name, kw)
            r = await self.execute(**kw)
            r.execution_time = time.time() - start
            self.count += 1
            viz.result(r.success, r.message or str(r.data)[:60] if r.data else "OK")
            if r.files_created: 
                for f in r.files_created: viz.file_created(f)
            return r
        except Exception as e: return ToolResult(False, self.name, error=str(e))


class ShellTool(BaseTool):
    name = "shell"; description = "Ejecuta comandos shell"
    def __init__(self, executor=None): super().__init__(); self.ex = executor or CommandExecutor()
    async def execute(self, command, timeout=30):
        r = await self.ex.execute(command, timeout)
        return ToolResult(r.success, self.name, data={'stdout': r.stdout, 'stderr': r.stderr}, 
                         message=r.stdout[:100] if r.stdout else "", error=r.error_message if not r.success else None)


class FileTool(BaseTool):
    name = "file"; description = "Operaciones de archivos"
    category = ToolCategory.FILE
    def __init__(self, fm=None): super().__init__(); self.fm = fm or FileManager()
    async def execute(self, operation, path, content=None, **kw):
        ops = {'read': lambda: self.fm.read(path), 'write': lambda: self.fm.write(path, content or ""),
               'delete': lambda: self.fm.delete(path), 'list': lambda: self.fm.list_dir(path),
               'mkdir': lambda: self.fm.mkdir(path), 'exists': lambda: self._exists(path)}
        if operation not in ops: return ToolResult(False, self.name, error=f"Op inválida: {operation}")
        r = await ops[operation]()
        if isinstance(r, bool): return ToolResult(True, self.name, data={'exists': r})
        fc = [path] if operation == 'write' and r.success else []
        return ToolResult(r.success, self.name, data=r.data, message=r.message, error=r.error, files_created=fc)
    async def _exists(self, p): return await self.fm.exists(p)


class PythonTool(BaseTool):
    name = "python"; description = "Ejecuta código Python"
    category = ToolCategory.DEVELOPMENT
    def __init__(self, ex=None): super().__init__(); self.ex = ex or CommandExecutor()
    async def execute(self, code, timeout=60):
        imports = "import sys,os,json,math,random,datetime\n"
        r = await self.ex.run_script(imports + code, 'python3', timeout)
        return ToolResult(r.success, self.name, data={'output': r.stdout, 'errors': r.stderr},
                         message=r.stdout[:150] if r.stdout else "", error=r.error_message if not r.success else None)


class SearchTool(BaseTool):
    name = "search"; description = "Búsqueda web en tiempo real"
    category = ToolCategory.SEARCH
    def __init__(self): super().__init__(); self.engine = WebSearchEngine()
    async def execute(self, query, num_results=10):
        r = await self.engine.search(query, num_results)
        viz.research(query, len(r.get('results', [])))
        return ToolResult(len(r.get('results', [])) > 0, self.name, data=r,
                         message=f"{len(r.get('results', []))} resultados para '{query}'")


class BrowserTool(BaseTool):
    name = "browser"; description = "Navega y extrae contenido web"
    category = ToolCategory.BROWSER
    def __init__(self): super().__init__(); self.browser = WebBrowser()
    async def execute(self, url):
        r = await self.browser.fetch(url)
        if r.get('error'): return ToolResult(False, self.name, error=r['error'])
        return ToolResult(True, self.name, data=r, message=f"Cargado: {r.get('title', url)[:50]}")


class DocumentTool(BaseTool):
    name = "document"; description = "Crea PowerPoint, Word, Excel"
    category = ToolCategory.DOCUMENT
    def __init__(self, out=None): super().__init__(); self.creator = DocumentCreator(out)
    async def execute(self, doc_type, title, content, filename=None, **kw):
        doc_type = doc_type.lower()
        if doc_type == 'pptx':
            slides = content if isinstance(content, list) else [{'title': 'Contenido', 'content': str(content)}]
            return await self.creator.create_pptx(title, slides, kw.get('theme', 'professional'), filename)
        elif doc_type == 'docx':
            return await self.creator.create_docx(title, content, kw.get('author', 'Agente IA'), filename)
        elif doc_type == 'xlsx':
            sheets = content if isinstance(content, list) else [{'name': 'Datos', 'headers': [], 'rows': []}]
            return await self.creator.create_xlsx(title, sheets, filename)
        return ToolResult(False, self.name, error=f"Tipo no soportado: {doc_type}")


class MessageTool(BaseTool):
    name = "message"; description = "Muestra mensajes"
    category = ToolCategory.COMMUNICATION
    async def execute(self, content, message_type="info", title=None):
        icons = {'info': 'ℹ️', 'success': '✅', 'warning': '⚠️', 'error': '❌', 'tip': '💡'}
        icon = icons.get(message_type, '💬')
        if title: console.print(Panel(content, title=f"{icon} {title}", border_style="cyan"))
        else: console.print(f"    {icon} {content}")
        return ToolResult(True, self.name, message=content)


class ResearchTool(BaseTool):
    name = "research"; description = "Investigación profunda sobre un tema"
    category = ToolCategory.SEARCH
    def __init__(self): super().__init__(); self.search = WebSearchEngine(); self.browser = WebBrowser()
    async def execute(self, topic, depth="basic", max_sources=5):
        viz.log(f"Investigando: {topic}", "🔬")
        sr = await self.search.search(topic, max_sources * 2)
        sources = []
        if depth in ["medium", "deep"] and sr.get('results'):
            for r in sr['results'][:max_sources]:
                if r.get('url'):
                    page = await self.browser.fetch(r['url'])
                    if not page.get('error'):
                        sources.append({'url': r['url'], 'title': page.get('title', ''), 'preview': page.get('content', '')[:400]})
        else:
            sources = sr.get('results', [])[:max_sources]
        
        data = {'topic': topic, 'sources': sources, 'summary': f"{len(sources)} fuentes sobre '{topic}'"}
        viz.research(topic, len(sources))
        return ToolResult(len(sources) > 0, self.name, data=data, message=f"Investigación: {len(sources)} fuentes")


class ToolRegistry:
    def __init__(self): self._tools: Dict[str, BaseTool] = {}
    def register(self, tool): self._tools[tool.name] = tool
    def get(self, name): return self._tools.get(name)
    async def execute(self, name, **params):
        tool = self.get(name)
        if not tool: return ToolResult(False, name, error=f"No encontrada: {name}")
        return await tool(**params)
    def list_tools(self): return list(self._tools.keys())
    def get_schemas(self): return [{'name': t.name, 'description': t.description, 'category': t.category.value} for t in self._tools.values()]


# ═══════════════════════════════════════════════════════════════════════════════
# PLANIFICADOR
# ═══════════════════════════════════════════════════════════════════════════════

class TaskPlanner:
    PATTERNS = {
        'create_pptx': [r'crea.*(?:pptx?|powerpoint|presentaci)', r'make.*(?:ppt|presentation|slides)'],
        'create_docx': [r'crea.*(?:docx?|word|documento)', r'make.*(?:doc|word)', r'escrib.*(?:informe|reporte)'],
        'create_xlsx': [r'crea.*(?:xlsx?|excel|hoja)', r'make.*(?:excel|spreadsheet)'],
        'research': [r'investig', r'research', r'analiz.*(?:tema|info)', r'busca.*informaci[oó]n.*sobre'],
        'search': [r'busc[ao]?\b', r'search\b', r'encuentra', r'find'],
        'browse': [r'naveg', r'visit', r'abre.*(?:url|p[aá]gina)', r'browse', r'fetch'],
        'file_create': [r'crea.*archivo', r'escrib.*archivo', r'make.*file'],
        'file_read': [r'lee.*archivo', r'muestra.*contenido', r'read.*file', r'cat\s'],
        'file_list': [r'lista.*archivos', r'ls\b', r'dir\b', r'list.*files'],
        'execute_code': [r'ejecut.*(?:c[oó]digo|python)', r'run.*(?:code|python)', r'corr[eo]'],
        'install': [r'instal', r'pip\s+install'],
        'help': [r'ayuda', r'help', r'c[oó]mo', r'qu[eé]\s+puedes'],
        'system': [r'sistema', r'status', r'info'],
    }
    
    def __init__(self):
        self._compiled = {k: [re.compile(p, re.I) for p in v] for k, v in self.PATTERNS.items()}
    
    def detect_intent(self, text):
        for intent, patterns in self._compiled.items():
            for p in patterns:
                if p.search(text): return intent
        return 'general'
    
    def extract_entities(self, text):
        return {'files': re.findall(r'[\w\-./]+\.\w{2,5}', text),
                'urls': re.findall(r'https?://[^\s<>"]+', text),
                'code': re.findall(r'```[\w]*\n?(.*?)```', text, re.DOTALL),
                'quoted': re.findall(r'["\']([^"\']+)["\']', text),
                'topic': (re.search(r'(?:sobre|about|de)\s+["\']?([^"\',.]+)', text, re.I) or type('',(),{'group':lambda s,x:None})()).group(1)}
    
    async def create_plan(self, user_input, context=None):
        intent = self.detect_intent(user_input)
        entities = self.extract_entities(user_input)
        viz.log(f"Intención: {intent}", "🎯")
        
        creators = {
            'create_pptx': self._pptx, 'create_docx': self._docx, 'create_xlsx': self._xlsx,
            'research': self._research, 'search': self._search, 'browse': self._browse,
            'file_create': self._file_create, 'file_read': self._file_read, 'file_list': self._file_list,
            'execute_code': self._code, 'install': self._install, 'help': self._help, 'system': self._system
        }
        return await creators.get(intent, self._general)(user_input, entities)
    
    async def _pptx(self, text, e):
        topic = e.get('topic') or (e['quoted'][0] if e.get('quoted') else "Presentación")
        return TaskPlan(uuid.uuid4().hex[:8], f"Crear PowerPoint: {topic}", [
            Phase("research", "Investigación", "Recopilar info", "🔬", [
                Step("s1", f"Investigar {topic}", "research", {"topic": topic, "depth": "basic"})]),
            Phase("create", "Creación", "Crear presentación", "📊", [
                Step("s2", "Generar slides", "document", {"doc_type": "pptx", "title": topic,
                    "content": [{"title": "Introducción", "content": f"Sobre {topic}"},
                               {"title": "Puntos Principales", "bullets": ["Punto 1", "Punto 2", "Punto 3"]},
                               {"title": "Conclusión", "content": "Resumen final"}]})]),
            Phase("deliver", "Entrega", "Confirmar", "✅", [
                Step("s3", "Confirmar", "message", {"content": f"Presentación '{topic}' creada", "message_type": "success"})])])
    
    async def _docx(self, text, e):
        topic = e.get('topic') or (e['quoted'][0] if e.get('quoted') else "Documento")
        return TaskPlan(uuid.uuid4().hex[:8], f"Crear Word: {topic}", [
            Phase("create", "Creación", "Crear documento", "📝", [
                Step("s1", "Generar documento", "document", {"doc_type": "docx", "title": topic,
                    "content": [{"title": "Introducción", "level": 1, "content": f"Este documento trata sobre {topic}."},
                               {"title": "Desarrollo", "level": 1, "content": "Contenido principal."},
                               {"title": "Conclusión", "level": 1, "content": "Resumen y conclusiones."}]})]),
            Phase("deliver", "Entrega", "Confirmar", "✅", [
                Step("s2", "Confirmar", "message", {"content": f"Documento '{topic}' creado", "message_type": "success"})])])
    
    async def _xlsx(self, text, e):
        topic = e.get('topic') or "Datos"
        return TaskPlan(uuid.uuid4().hex[:8], f"Crear Excel: {topic}", [
            Phase("create", "Creación", "Crear Excel", "📈", [
                Step("s1", "Generar Excel", "document", {"doc_type": "xlsx", "title": topic,
                    "content": [{"name": "Datos", "headers": ["Columna A", "Columna B", "Columna C"],
                                "rows": [["Dato 1", 100, 50], ["Dato 2", 200, 75], ["Dato 3", 150, 60]],
                                "chart": {"type": "bar", "title": "Gráfico"}}]})]),
            Phase("deliver", "Entrega", "Confirmar", "✅", [
                Step("s2", "Confirmar", "message", {"content": f"Excel '{topic}' creado", "message_type": "success"})])])
    
    async def _research(self, text, e):
        topic = e.get('topic') or text.replace('investiga', '').replace('sobre', '').strip()[:50]
        return TaskPlan(uuid.uuid4().hex[:8], f"Investigar: {topic}", [
            Phase("search", "Búsqueda", "Buscar info", "🔍", [
                Step("s1", f"Buscar sobre {topic}", "search", {"query": topic})]),
            Phase("analyze", "Análisis", "Analizar fuentes", "🔬", [
                Step("s2", "Investigación", "research", {"topic": topic, "depth": "medium", "max_sources": 5})]),
            Phase("report", "Reporte", "Mostrar resultados", "📋", [
                Step("s3", "Mostrar", "message", {"content": f"Investigación sobre '{topic}' completada", "message_type": "success"})])])
    
    async def _search(self, text, e):
        q = text.replace('busca', '').replace('buscar', '').strip()[:100]
        return TaskPlan(uuid.uuid4().hex[:8], f"Buscar: {q}", [
            Phase("search", "Búsqueda", "Buscar", "🔍", [Step("s1", f"Buscar: {q}", "search", {"query": q})])])
    
    async def _browse(self, text, e):
        url = e['urls'][0] if e.get('urls') else "https://example.com"
        return TaskPlan(uuid.uuid4().hex[:8], f"Navegar: {url[:30]}", [
            Phase("browse", "Navegación", "Obtener contenido", "🌐", [Step("s1", f"Navegar", "browser", {"url": url})])])
    
    async def _file_create(self, text, e):
        fn = e['files'][0] if e.get('files') else 'nuevo_archivo.txt'
        content = e['quoted'][0] if e.get('quoted') else "# Nuevo archivo\n"
        return TaskPlan(uuid.uuid4().hex[:8], f"Crear: {fn}", [
            Phase("create", "Crear", "Crear archivo", "📄", [Step("s1", f"Crear {fn}", "file", {"operation": "write", "path": fn, "content": content})])])
    
    async def _file_read(self, text, e):
        fn = e['files'][0] if e.get('files') else ''
        return TaskPlan(uuid.uuid4().hex[:8], f"Leer: {fn}", [
            Phase("read", "Leer", "Leer archivo", "📖", [Step("s1", f"Leer {fn}", "file", {"operation": "read", "path": fn})])])
    
    async def _file_list(self, text, e):
        return TaskPlan(uuid.uuid4().hex[:8], "Listar archivos", [
            Phase("list", "Listar", "Mostrar archivos", "📁", [Step("s1", "Listar", "file", {"operation": "list", "path": "."})])])
    
    async def _code(self, text, e):
        code = e['code'][0] if e.get('code') else 'print("Hola Mundo!")'
        return TaskPlan(uuid.uuid4().hex[:8], "Ejecutar Python", [
            Phase("exec", "Ejecutar", "Ejecutar código", "🐍", [Step("s1", "Ejecutar", "python", {"code": code})])])
    
    async def _install(self, text, e):
        m = re.search(r'install\s+(\S+)', text, re.I)
        pkg = m.group(1) if m else 'requests'
        return TaskPlan(uuid.uuid4().hex[:8], f"Instalar: {pkg}", [
            Phase("install", "Instalar", f"pip install {pkg}", "📦", [Step("s1", f"Instalar {pkg}", "shell", {"command": f"pip install {pkg} -q"})])])
    
    async def _help(self, text, e):
        help_msg = """🤖 **Agente IA v2.0** - Capacidades:

📊 **Documentos**: Crear PowerPoint, Word, Excel
🔍 **Investigación**: Buscar y analizar información web
🌐 **Navegación**: Visitar y extraer contenido de URLs
📁 **Archivos**: Crear, leer, listar, copiar, mover
🐍 **Código**: Ejecutar Python
📦 **Paquetes**: Instalar dependencias

**Ejemplos:**
• "Crea una presentación sobre inteligencia artificial"
• "Investiga sobre cambio climático"
• "Crea un documento Word con un informe"
• "Busca noticias sobre tecnología"
• "Lista los archivos del directorio" """
        return TaskPlan(uuid.uuid4().hex[:8], "Ayuda", [
            Phase("help", "Ayuda", "Mostrar info", "❓", [Step("s1", "Mostrar", "message", {"content": help_msg, "message_type": "info", "title": "Ayuda"})])])
    
    async def _system(self, text, e):
        return TaskPlan(uuid.uuid4().hex[:8], "Info del sistema", [
            Phase("info", "Sistema", "Obtener info", "💻", [Step("s1", "Info", "shell", {"command": "uname -a && python3 --version && pwd"})])])
    
    async def _general(self, text, e):
        return TaskPlan(uuid.uuid4().hex[:8], text[:80], [
            Phase("process", "Procesar", "Procesar solicitud", "🔄", [Step("s1", "Analizar", "message", {"content": f"Procesando: {text[:80]}...", "message_type": "info"})])])


# ═══════════════════════════════════════════════════════════════════════════════
# AGENTE PRINCIPAL v2.0
# ═══════════════════════════════════════════════════════════════════════════════

@dataclass
class AgentConfig:
    name: str = "Agente IA v2.0"
    max_iterations: int = 100
    timeout: int = 60
    verbose: bool = True

class Agent:
    """Agente IA v2.0 con ciclo de 7 pasos y visualización en tiempo real."""
    
    def __init__(self, config=None):
        self.config = config or AgentConfig()
        self.state = AgentState.IDLE
        
        self.security = SecurityGuard()
        self.executor = CommandExecutor(self.security, self.config.timeout)
        self.files = FileManager(self.security)
        
        self.tools = ToolRegistry()
        self._register_tools()
        self.planner = TaskPlanner()
        
        self.current_plan = None
        self.history = []
        self.iteration = 0
        
        viz.start()
    
    def _register_tools(self):
        self.tools.register(ShellTool(self.executor))
        self.tools.register(FileTool(self.files))
        self.tools.register(PythonTool(self.executor))
        self.tools.register(SearchTool())
        self.tools.register(BrowserTool())
        self.tools.register(DocumentTool(DEFAULT_OUTPUT))
        self.tools.register(MessageTool())
        self.tools.register(ResearchTool())
    
    async def run(self, user_input):
        self.iteration = 0
        self.state = AgentState.ANALYZING
        
        try:
            if self.config.verbose:
                console.print()
                console.print(Rule(f"[bold blue]🤖 {self.config.name}[/bold blue]"))
                console.print(f"[dim]📝 {user_input[:80]}{'...' if len(user_input) > 80 else ''}[/dim]\n")
            
            # Planificar
            viz.phase("Planificación", "📋")
            self.state = AgentState.PLANNING
            self.current_plan = await self.planner.create_plan(user_input)
            
            if self.current_plan.phases:
                self.current_plan.phases[0].status = PhaseStatus.IN_PROGRESS
            
            if self.config.verbose:
                self._show_plan()
            
            # Ejecutar
            while not self.current_plan.is_complete and self.iteration < self.config.max_iterations:
                self.iteration += 1
                phase = self.current_plan.get_current_phase()
                if not phase: break
                
                viz.phase(phase.name, phase.icon)
                
                step = phase.get_next_step()
                if not step:
                    self.current_plan.advance()
                    continue
                
                self.state = AgentState.EXECUTING
                viz.step(step.description)
                result = await self.tools.execute(step.tool, **step.params)
                
                if result.success:
                    step.complete(result.data)
                else:
                    step.fail(result.error or "Error")
                
                self.history.append({'tool': step.tool, 'params': step.params, 'success': result.success})
                self.current_plan.advance()
            
            # Entregar
            self.state = AgentState.DELIVERING
            response = self._build_response()
            
            if self.config.verbose:
                viz.summary(self.current_plan)
            
            return response
        
        except Exception as e:
            self.state = AgentState.ERROR
            console.print(f"[red]❌ Error: {e}[/red]")
            return f"Error: {e}"
        finally:
            self.state = AgentState.IDLE
    
    def _show_plan(self):
        tree = Tree(f"[bold]📋 {self.current_plan.objective[:50]}[/bold]")
        for p in self.current_plan.phases:
            branch = tree.add(f"{p.icon} [cyan]{p.name}[/cyan]")
            for s in p.steps: branch.add(f"[dim]• {s.description}[/dim]")
        console.print(tree); console.print()
    
    def _build_response(self):
        if not self.current_plan: return "Sin plan."
        parts = [f"**{self.current_plan.objective}**\n"]
        files, results = [], []
        for p in self.current_plan.phases:
            for s in p.steps:
                if s.result:
                    if isinstance(s.result, dict):
                        if s.result.get('path'): files.append(s.result['path'])
                        if s.result.get('stdout'): results.append(s.result['stdout'][:150])
                        if s.result.get('output'): results.append(s.result['output'][:150])
                        if s.result.get('results'):
                            for r in s.result['results'][:3]:
                                results.append(f"• {r.get('title', '')[:60]}")
                    else: results.append(str(s.result)[:150])
        if files:
            parts.append("\n📁 **Archivos:**")
            for f in files: parts.append(f"  • {f}")
        if results:
            parts.append("\n📊 **Resultados:**")
            for r in results[:5]: parts.append(f"  {r}")
        parts.append(f"\n✅ Completado en {self.iteration} pasos")
        return "\n".join(parts)
    
    async def execute_direct(self, tool_name, **params):
        return await self.tools.execute(tool_name, **params)
    
    def get_status(self):
        return {'name': self.config.name, 'state': self.state.value, 'iterations': self.iteration,
                'tools': len(self.tools.list_tools()), 'history': len(self.history)}

# ═══════════════════════════════════════════════════════════════════════════════
# INTERFAZ DE USUARIO
# ═══════════════════════════════════════════════════════════════════════════════

def print_banner():
    console.print("""
[bold blue]
╔══════════════════════════════════════════════════════════════════════════════╗
║            🚀 AGENTE IA v2.0 - SISTEMA COMPLETO 🚀                          ║
╠══════════════════════════════════════════════════════════════════════════════╣
║  ✨ PowerPoint, Word, Excel  |  🔍 Búsqueda Web  |  🌐 Navegación           ║
║  📁 Gestión de Archivos      |  🐍 Python       |  🔬 Investigación         ║
║  👁️ Visualización en Tiempo Real del Proceso Completo                       ║
╚══════════════════════════════════════════════════════════════════════════════╝
[/bold blue]""")

def print_help():
    console.print(Panel("""
[bold cyan]📚 COMANDOS[/bold cyan]

[yellow]Sistema:[/yellow]
  exit, salir     Salir
  help, ayuda     Esta ayuda
  status          Estado del agente
  tools           Herramientas
  clear           Limpiar

[yellow]Directo:[/yellow]
  !<cmd>          Shell (ej: !ls -la)
  py:<código>     Python (ej: py:print(2+2))

[yellow]Documentos:[/yellow]
  "Crea una presentación sobre [tema]"
  "Crea un documento Word sobre [tema]"
  "Crea un Excel con [datos]"

[yellow]Investigación:[/yellow]
  "Investiga sobre [tema]"
  "Busca información de [query]"
  "Navega a [url]"

[yellow]Archivos:[/yellow]
  "Lista los archivos"
  "Crea un archivo [nombre]"
""", title="[bold]Ayuda[/bold]", border_style="cyan"))


async def demo():
    print_banner()
    console.print("\n[bold cyan]🎯 DEMO DEL AGENTE IA v2.0[/bold cyan]\n")
    agent = Agent(AgentConfig(verbose=True))
    demos = [("Sistema", "muestra información del sistema"), ("Archivo", "crea archivo demo.txt con 'Hola v2.0!'"),
             ("Listar", "lista los archivos"), ("Python", "ejecuta: print([x**2 for x in range(5)])")]
    for i, (title, task) in enumerate(demos, 1):
        console.print(f"\n[yellow]═══ Demo {i}: {title} ═══[/yellow]")
        result = await agent.run(task)
        console.print(f"\n[dim]{result[:200]}[/dim]")
        await asyncio.sleep(0.5)
    console.print("\n[green]✅ Demo completada[/green]")


async def interactive():
    print_banner()
    agent = Agent(AgentConfig(verbose=True))
    console.print("\n[green]✅ Agente listo[/green] - 'help' para ayuda\n")
    
    while True:
        try:
            inp = Prompt.ask("[bold green]🤖 v2.0[/bold green]").strip()
            if not inp: continue
            cmd = inp.lower()
            
            if cmd in ['exit', 'salir', 'quit', 'q']:
                console.print("\n[yellow]👋 ¡Hasta luego![/yellow]"); break
            elif cmd in ['help', 'ayuda', '?']:
                print_help(); continue
            elif cmd == 'status':
                s = agent.get_status()
                t = Table(show_header=False, box=box.SIMPLE)
                t.add_column("", style="cyan"); t.add_column("", style="green")
                for k, v in s.items(): t.add_row(str(k), str(v))
                console.print(t); continue
            elif cmd == 'tools':
                t = Table(title="Herramientas", box=box.ROUNDED)
                t.add_column("Nombre", style="cyan"); t.add_column("Descripción")
                for tool in agent.tools.get_schemas(): t.add_row(tool['name'], tool['description'][:40])
                console.print(t); continue
            elif cmd == 'clear':
                console.clear(); print_banner(); continue
            elif inp.startswith('!'):
                r = await agent.execute_direct('shell', command=inp[1:].strip())
                if r.success and r.data: console.print(r.data.get('stdout', ''))
                else: console.print(f"[red]{r.error}[/red]")
                continue
            elif inp.startswith('py:'):
                r = await agent.execute_direct('python', code=inp[3:].strip())
                if r.success and r.data: console.print(r.data.get('output', ''))
                else: console.print(f"[red]{r.error}[/red]")
                continue
            
            result = await agent.run(inp)
            console.print(f"\n[bold]📋 Resultado:[/bold]\n{result}\n")
        
        except KeyboardInterrupt:
            console.print("\n[yellow]Ctrl+C para salir[/yellow]")
            try: await asyncio.sleep(1)
            except KeyboardInterrupt: console.print("\n[yellow]👋 ¡Adiós![/yellow]"); break
        except Exception as e: console.print(f"[red]Error: {e}[/red]")


async def tests():
    console.print("\n[bold cyan]🧪 TESTS[/bold cyan]\n")
    passed = 0
    
    console.print("Test 1: Security...", end=" ")
    s = SecurityGuard()
    if s.analyze("ls").is_safe and not s.analyze("rm -rf /").is_safe: console.print("[green]✅[/green]"); passed += 1
    else: console.print("[red]❌[/red]")
    
    console.print("Test 2: Executor...", end=" ")
    e = CommandExecutor(s)
    r = await e.execute("echo test")
    if r.success and "test" in r.stdout: console.print("[green]✅[/green]"); passed += 1
    else: console.print("[red]❌[/red]")
    
    console.print("Test 3: FileManager...", end=" ")
    fm = FileManager(s)
    await fm.write("test.txt", "contenido")
    r = await fm.read("test.txt")
    if r.success and r.data == "contenido": console.print("[green]✅[/green]"); passed += 1
    else: console.print("[red]❌[/red]")
    await fm.delete("test.txt")
    
    console.print("Test 4: Agent...", end=" ")
    a = Agent(AgentConfig(verbose=False))
    r = await a.run("ayuda")
    if r and len(r) > 0: console.print("[green]✅[/green]"); passed += 1
    else: console.print("[red]❌[/red]")
    
    console.print(f"\n[bold]{passed}/4 tests pasados[/bold]")
    if passed == 4: console.print("[green]✅ Todos OK[/green]")


async def main():
    print_banner()
    console.print("\n[bold]¿Qué hacer?[/bold]")
    console.print("  [cyan]1.[/cyan] Modo interactivo")
    console.print("  [cyan]2.[/cyan] Demo rápida")
    console.print("  [cyan]3.[/cyan] Tests")
    console.print("  [cyan]0.[/cyan] Salir")
    
    choice = Prompt.ask("\n[bold]Opción[/bold]", default="1")
    if choice == "1": await interactive()
    elif choice == "2": await demo()
    elif choice == "3": await tests()
    else: console.print("[yellow]👋 ¡Adiós![/yellow]")

if __name__ == "__main__":
    try: asyncio.run(main())
    except KeyboardInterrupt: console.print("\n[yellow]👋 ¡Adiós![/yellow]")
    except Exception as e: console.print(f"[red]Error: {e}[/red]"); traceback.print_exc()
