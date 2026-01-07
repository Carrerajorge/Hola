#!/usr/bin/env python3
"""
╔═════════════════════════════════════════════════════════════════════════════════════════════════════╗
║                                                                                                     ║
║     █████╗  ██████╗ ███████╗███╗   ██╗████████╗███████╗    ██╗ █████╗    ██╗   ██╗███████╗         ║
║    ██╔══██╗██╔════╝ ██╔════╝████╗  ██║╚══██╔══╝██╔════╝    ██║██╔══██╗   ██║   ██║██╔════╝         ║
║    ███████║██║  ███╗█████╗  ██╔██╗ ██║   ██║   █████╗      ██║███████║   ██║   ██║███████╗         ║
║    ██╔══██║██║   ██║██╔══╝  ██║╚██╗██║   ██║   ██╔══╝      ██║██╔══██║   ╚██╗ ██╔╝╚════██║         ║
║    ██║  ██║╚██████╔╝███████╗██║ ╚████║   ██║   ███████╗    ██║██║  ██║    ╚████╔╝ ███████║         ║
║    ╚═╝  ╚═╝ ╚═════╝ ╚══════╝╚═╝  ╚═══╝   ╚═╝   ╚══════╝    ╚═╝╚═╝  ╚═╝     ╚═══╝  ╚══════╝         ║
║                                                                                                     ║
║                     🚀 SISTEMA AGENTE IA v5.0 - ENTERPRISE EDITION 🚀                               ║
║                                                                                                     ║
╠═════════════════════════════════════════════════════════════════════════════════════════════════════╣
║                                                                                                     ║
║  🌐 NAVEGADOR WEB INTEGRADO (Playwright/Selenium)                                                   ║
║  ├── 🖥️ Computadora Virtual - Navega como humano                                                    ║
║  ├── 📸 Screenshots y grabación de pantalla                                                         ║
║  ├── 🔄 Pool de navegadores reutilizables                                                           ║
║  ├── 🛡️ Anti-detección y rotación de user agents                                                   ║
║  └── ⚡ Caché inteligente para rendimiento                                                          ║
║                                                                                                     ║
║  📊 DOCUMENTOS PROFESIONALES                                                                        ║
║  ├── 📑 PowerPoint con diseños premium                                                              ║
║  ├── 📝 Word con formato avanzado                                                                   ║
║  ├── 📈 Excel con fórmulas y gráficos                                                               ║
║  └── 📄 PDF generation                                                                              ║
║                                                                                                     ║
║  🔧 ARQUITECTURA ENTERPRISE                                                                         ║
║  ├── 📬 Sistema de colas (Queue) para tareas                                                        ║
║  ├── 💾 Caché multinivel (Memory + Disk)                                                            ║
║  ├── 🔒 Rate limiting por dominio y usuario                                                         ║
║  ├── 📊 Métricas y monitoreo en tiempo real                                                         ║
║  ├── 🔄 Workers concurrentes                                                                        ║
║  └── 🛡️ Sistema de seguridad multicapa                                                             ║
║                                                                                                     ║
║  👁️ VISUALIZACIÓN EN TIEMPO REAL DE TODO EL PROCESO                                                ║
║                                                                                                     ║
╚═════════════════════════════════════════════════════════════════════════════════════════════════════╝

INSTRUCCIONES PARA REPLIT:
1. Crea un nuevo Repl Python 3.10+
2. Copia este archivo como main.py
3. Ejecuta - instalará TODAS las dependencias
4. ¡Disfruta el poder del Agente IA v5.0!
"""

from __future__ import annotations
import os, sys, re, json, shutil, hashlib, asyncio, tempfile, subprocess, shlex
import mimetypes, time, logging, uuid, traceback, base64, io, threading
import pickle, gzip, sqlite3, queue as queue_module, weakref, functools
from pathlib import Path
from datetime import datetime, timedelta
from typing import Dict, Any, Optional, List, Union, Callable, Set, Tuple
from dataclasses import dataclass, field, asdict
from enum import Enum, auto
from abc import ABC, abstractmethod
from collections import deque, defaultdict, OrderedDict
from concurrent.futures import ThreadPoolExecutor, as_completed
from contextlib import asynccontextmanager, contextmanager
from urllib.parse import urlparse, urljoin, quote_plus
import inspect
import struct
import socket
import random
import fnmatch

# ═══════════════════════════════════════════════════════════════════════════════
# INSTALACIÓN AVANZADA DE DEPENDENCIAS
# ═══════════════════════════════════════════════════════════════════════════════

class DependencyInstaller:
    """Instalador inteligente de dependencias con reintentos."""
    
    CORE_PACKAGES = [
        ("aiofiles", "aiofiles"),
        ("rich", "rich"),
        ("httpx", "httpx"),
        ("aiohttp", "aiohttp"),
        ("beautifulsoup4", "bs4"),
        ("lxml", "lxml"),
        ("python-pptx", "pptx"),
        ("python-docx", "docx"),
        ("openpyxl", "openpyxl"),
        ("Pillow", "PIL"),
        ("fake-useragent", "fake_useragent"),
        ("diskcache", "diskcache"),
    ]
    
    BROWSER_PACKAGES = [
        ("playwright", "playwright"),
        ("selenium", "selenium"),
        ("webdriver-manager", "webdriver_manager"),
    ]
    
    OPTIONAL_PACKAGES = [
        ("pandas", "pandas"),
        ("matplotlib", "matplotlib"),
        ("numpy", "numpy"),
    ]
    
    @classmethod
    def install(cls, pkg, retries=2):
        # Validate package name to prevent command injection
        # Allow: letters, numbers, hyphens, underscores, periods, brackets for extras
        if not re.match(r'^[a-zA-Z0-9._\-\[\],<>=!]+$', pkg):
            print(f"⚠️ Invalid package name: {pkg}")
            return False
        
        for i in range(retries):
            try:
                subprocess.run([sys.executable, "-m", "pip", "install", pkg, "-q", "--disable-pip-version-check"],
                              capture_output=True, timeout=120)
                return True
            except: pass
        return False
    
    @classmethod
    def setup(cls):
        print("\n🔧 Configurando Agente IA v5.0...")
        print("═" * 60)
        
        # Core packages
        for pip_name, import_name in cls.CORE_PACKAGES:
            try:
                __import__(import_name)
                print(f"  ✅ {pip_name}")
            except ImportError:
                print(f"  📦 Instalando {pip_name}...", end=" ", flush=True)
                if cls.install(pip_name):
                    try: __import__(import_name); print("✅")
                    except: print("⚠️")
                else: print("❌")
        
        # Browser packages (opcional pero recomendado)
        print("\n  🌐 Navegador web:")
        for pip_name, import_name in cls.BROWSER_PACKAGES:
            try:
                __import__(import_name)
                print(f"    ✅ {pip_name}")
            except ImportError:
                print(f"    📦 {pip_name}...", end=" ", flush=True)
                if cls.install(pip_name):
                    try: __import__(import_name); print("✅")
                    except: print("⚠️")
                else: print("⚠️ (opcional)")
        
        # Playwright browsers
        try:
            import playwright
            print("    📦 Instalando navegadores Playwright...", end=" ", flush=True)
            subprocess.run([sys.executable, "-m", "playwright", "install", "chromium"],
                          capture_output=True, timeout=300)
            print("✅")
        except: print("    ⚠️ Navegadores Playwright (se instalará al usar)")
        
        print("═" * 60)
        print("✨ Configuración completada\n")

# Ejecutar instalación
DependencyInstaller.setup()

# ═══════════════════════════════════════════════════════════════════════════════
# IMPORTACIONES POST-INSTALACIÓN
# ═══════════════════════════════════════════════════════════════════════════════

import aiofiles
from rich.console import Console, Group
from rich.panel import Panel
from rich.table import Table
from rich.tree import Tree
from rich.rule import Rule
from rich.prompt import Prompt, Confirm
from rich.progress import Progress, SpinnerColumn, TextColumn, BarColumn, TaskProgressColumn
from rich.live import Live
from rich.layout import Layout
from rich.syntax import Syntax
from rich.markdown import Markdown
from rich.text import Text
from rich import box

console = Console()
logging.basicConfig(level=logging.WARNING, format='%(message)s')
logger = logging.getLogger("AgentIA_v5")

# Importaciones opcionales
try: import httpx; HTTPX_OK = True
except: HTTPX_OK = False

try: import aiohttp; AIOHTTP_OK = True
except: AIOHTTP_OK = False

try: from bs4 import BeautifulSoup; BS4_OK = True
except: BS4_OK = False

try: from lxml import html as lxml_html; LXML_OK = True
except: LXML_OK = False

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    PPTX_OK = True
except: PPTX_OK = False

try:
    from docx import Document as DocxDocument
    from docx.shared import Inches as DI, Pt as DP, RGBColor as DRC
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.enum.table import WD_TABLE_ALIGNMENT
    DOCX_OK = True
except: DOCX_OK = False

try:
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side, NamedStyle
    from openpyxl.chart import BarChart, LineChart, PieChart, Reference
    from openpyxl.utils import get_column_letter
    from openpyxl.formatting.rule import ColorScaleRule, FormulaRule
    XLSX_OK = True
except: XLSX_OK = False

try: from PIL import Image, ImageDraw, ImageFont, ImageFilter; PIL_OK = True
except: PIL_OK = False

try: from fake_useragent import UserAgent; UA_OK = True
except: UA_OK = False

try: import diskcache; CACHE_OK = True
except: CACHE_OK = False

# Navegadores
PLAYWRIGHT_OK = False
SELENIUM_OK = False

try:
    from playwright.async_api import async_playwright, Browser, Page, BrowserContext
    PLAYWRIGHT_OK = True
except: pass

try:
    from selenium import webdriver
    from selenium.webdriver.chrome.options import Options as ChromeOptions
    from selenium.webdriver.common.by import By
    from selenium.webdriver.support.ui import WebDriverWait
    from selenium.webdriver.support import expected_conditions as EC
    from selenium.common.exceptions import TimeoutException, WebDriverException
    try:
        from webdriver_manager.chrome import ChromeDriverManager
        from selenium.webdriver.chrome.service import Service
    except: pass
    SELENIUM_OK = True
except: pass

# ═══════════════════════════════════════════════════════════════════════════════
# CONSTANTES Y CONFIGURACIÓN GLOBAL
# ═══════════════════════════════════════════════════════════════════════════════

VERSION = "5.0.0"
WORKSPACE = Path(os.getcwd()) / "workspace_v5"
OUTPUT_DIR = WORKSPACE / "output"
CACHE_DIR = WORKSPACE / ".cache"
SCREENSHOTS_DIR = WORKSPACE / "screenshots"
TEMP_DIR = WORKSPACE / ".temp"

# Crear directorios
for d in [WORKSPACE, OUTPUT_DIR, CACHE_DIR, SCREENSHOTS_DIR, TEMP_DIR]:
    d.mkdir(parents=True, exist_ok=True)

# Límites
MAX_FILE_SIZE = 500 * 1024 * 1024  # 500MB
MAX_TIMEOUT = 300  # 5 min
MAX_CONCURRENT_BROWSERS = 5
MAX_CACHE_SIZE = 1000
CACHE_TTL = 3600  # 1 hora

# User Agents
USER_AGENTS = [
    "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36",
    "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36",
    "Mozilla/5.0 (Windows NT 10.0; Win64; x64; rv:121.0) Gecko/20100101 Firefox/121.0",
    "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) AppleWebKit/605.1.15 (KHTML, like Gecko) Version/17.2 Safari/605.1.15",
]

def get_random_ua():
    if UA_OK:
        try: return UserAgent().random
        except: pass
    return random.choice(USER_AGENTS)


# ═══════════════════════════════════════════════════════════════════════════════
# ENUMERACIONES Y DATACLASSES
# ═══════════════════════════════════════════════════════════════════════════════

class ThreatLevel(Enum):
    SAFE = "safe"; LOW = "low"; MEDIUM = "medium"; HIGH = "high"; CRITICAL = "critical"

class SecurityAction(Enum):
    ALLOW = "allow"; WARN = "warn"; BLOCK = "block"; LOG_AND_BLOCK = "log_and_block"

class ExecutionStatus(Enum):
    PENDING = "pending"; QUEUED = "queued"; RUNNING = "running"; COMPLETED = "completed"
    FAILED = "failed"; TIMEOUT = "timeout"; BLOCKED = "blocked"; CACHED = "cached"

class AgentState(Enum):
    IDLE = "idle"; ANALYZING = "analyzing"; PLANNING = "planning"; EXECUTING = "executing"
    BROWSING = "browsing"; RESEARCHING = "researching"; CREATING = "creating"
    DELIVERING = "delivering"; ERROR = "error"

class PhaseStatus(Enum):
    PENDING = "pending"; IN_PROGRESS = "in_progress"; COMPLETED = "completed"; FAILED = "failed"

class ToolCategory(Enum):
    SYSTEM = "system"; FILE = "file"; DOCUMENT = "document"; SEARCH = "search"
    BROWSER = "browser"; SCRAPER = "scraper"; COMMUNICATION = "communication"
    DEVELOPMENT = "development"; DATA = "data"; AI = "ai"

class BrowserEngine(Enum):
    PLAYWRIGHT = "playwright"; SELENIUM = "selenium"; HTTPX = "httpx"

@dataclass
class SecurityAnalysis:
    command: str; is_safe: bool; threat_level: ThreatLevel
    action: SecurityAction; warnings: List[str] = field(default_factory=list)
    risk_score: float = 0.0

@dataclass
class ExecutionResult:
    command: str; status: ExecutionStatus; return_code: int = None
    stdout: str = ""; stderr: str = ""; execution_time: float = 0.0
    error_message: str = ""; cached: bool = False
    @property
    def success(self): return self.status in [ExecutionStatus.COMPLETED, ExecutionStatus.CACHED] and (self.return_code is None or self.return_code == 0)

@dataclass
class FileOperationResult:
    success: bool; operation: str; path: str
    message: str = ""; data: Any = None; error: str = None
    bytes_processed: int = 0

@dataclass
class ToolResult:
    success: bool; tool_name: str; data: Any = None
    message: str = ""; error: str = None; execution_time: float = 0.0
    files_created: List[str] = field(default_factory=list)
    screenshots: List[str] = field(default_factory=list)
    cached: bool = False

@dataclass
class WebPage:
    url: str; final_url: str = ""; title: str = ""; text: str = ""
    html: str = ""; links: List[str] = field(default_factory=list)
    images: List[str] = field(default_factory=list)
    metadata: Dict[str, str] = field(default_factory=dict)
    screenshot: str = None; status_code: int = 0
    load_time: float = 0.0; error: str = None
    js_rendered: bool = False

@dataclass
class SearchResult:
    query: str; results: List[Dict] = field(default_factory=list)
    total: int = 0; search_time: float = 0.0
    source: str = ""; cached: bool = False

@dataclass
class Step:
    id: str; description: str; tool: str; params: Dict[str, Any]
    status: PhaseStatus = PhaseStatus.PENDING; result: Any = None
    error: str = None; start_time: datetime = None; end_time: datetime = None
    def start(self): self.status = PhaseStatus.IN_PROGRESS; self.start_time = datetime.now()
    def complete(self, r): self.status = PhaseStatus.COMPLETED; self.result = r; self.end_time = datetime.now()
    def fail(self, e): self.status = PhaseStatus.FAILED; self.error = e; self.end_time = datetime.now()
    @property
    def duration(self): return (self.end_time - self.start_time).total_seconds() if self.start_time and self.end_time else 0

@dataclass
class Phase:
    id: str; name: str; description: str; icon: str = "📋"
    steps: List[Step] = field(default_factory=list)
    status: PhaseStatus = PhaseStatus.PENDING
    @property
    def pending_steps(self): return [s for s in self.steps if s.status == PhaseStatus.PENDING]
    @property
    def is_complete(self): return len(self.steps) > 0 and all(s.status == PhaseStatus.COMPLETED for s in self.steps)
    def get_next_step(self): return self.pending_steps[0] if self.pending_steps else None

@dataclass
class TaskPlan:
    task_id: str; objective: str
    phases: List[Phase] = field(default_factory=list)
    current_phase_index: int = 0
    created_at: datetime = field(default_factory=datetime.now)
    @property
    def is_complete(self): return all(p.is_complete for p in self.phases)
    @property
    def progress(self):
        total = sum(len(p.steps) for p in self.phases)
        done = sum(1 for p in self.phases for s in p.steps if s.status == PhaseStatus.COMPLETED)
        return (done / total) * 100 if total > 0 else 0
    def get_current_phase(self):
        return self.phases[self.current_phase_index] if 0 <= self.current_phase_index < len(self.phases) else None
    def advance(self):
        cur = self.get_current_phase()
        if cur and cur.is_complete:
            cur.status = PhaseStatus.COMPLETED
            self.current_phase_index += 1
            nxt = self.get_current_phase()
            if nxt: nxt.status = PhaseStatus.IN_PROGRESS

# ═══════════════════════════════════════════════════════════════════════════════
# SISTEMA DE CACHÉ MULTINIVEL
# ═══════════════════════════════════════════════════════════════════════════════

class CacheManager:
    """Caché multinivel: Memoria + Disco."""
    
    def __init__(self, max_memory=500, disk_dir=None):
        self.memory_cache: OrderedDict = OrderedDict()
        self.max_memory = max_memory
        self.disk_dir = Path(disk_dir or CACHE_DIR)
        self.disk_dir.mkdir(parents=True, exist_ok=True)
        self.stats = {'hits': 0, 'misses': 0, 'memory_hits': 0, 'disk_hits': 0}
        self._lock = threading.Lock()
        
        # DiskCache si está disponible
        self.disk_cache = None
        if CACHE_OK:
            try:
                self.disk_cache = diskcache.Cache(str(self.disk_dir / "diskcache"))
            except: pass
    
    def _hash_key(self, key: str) -> str:
        return hashlib.md5(key.encode()).hexdigest()
    
    def get(self, key: str) -> Optional[Any]:
        hkey = self._hash_key(key)
        
        # Memoria
        with self._lock:
            if hkey in self.memory_cache:
                value, expires = self.memory_cache[hkey]
                if expires is None or datetime.now() < expires:
                    self.memory_cache.move_to_end(hkey)
                    self.stats['hits'] += 1
                    self.stats['memory_hits'] += 1
                    return value
                else:
                    del self.memory_cache[hkey]
        
        # Disco
        if self.disk_cache:
            try:
                value = self.disk_cache.get(hkey)
                if value is not None:
                    self.stats['hits'] += 1
                    self.stats['disk_hits'] += 1
                    # Promover a memoria
                    self._set_memory(hkey, value, CACHE_TTL)
                    return value
            except: pass
        
        self.stats['misses'] += 1
        return None
    
    def set(self, key: str, value: Any, ttl: int = CACHE_TTL):
        hkey = self._hash_key(key)
        self._set_memory(hkey, value, ttl)
        
        # También en disco
        if self.disk_cache:
            try:
                self.disk_cache.set(hkey, value, expire=ttl)
            except: pass
    
    def _set_memory(self, hkey: str, value: Any, ttl: int):
        expires = datetime.now() + timedelta(seconds=ttl) if ttl else None
        with self._lock:
            if hkey in self.memory_cache:
                self.memory_cache.move_to_end(hkey)
            self.memory_cache[hkey] = (value, expires)
            
            # Evict si excede límite
            while len(self.memory_cache) > self.max_memory:
                self.memory_cache.popitem(last=False)
    
    def clear(self):
        with self._lock:
            self.memory_cache.clear()
        if self.disk_cache:
            try: self.disk_cache.clear()
            except: pass
    
    def get_stats(self) -> Dict:
        total = self.stats['hits'] + self.stats['misses']
        return {
            **self.stats,
            'hit_rate': (self.stats['hits'] / total * 100) if total > 0 else 0,
            'memory_size': len(self.memory_cache)
        }

# Cache global
cache = CacheManager()

# ═══════════════════════════════════════════════════════════════════════════════
# RATE LIMITER
# ═══════════════════════════════════════════════════════════════════════════════

class RateLimiter:
    """Rate limiting por dominio y global."""
    
    def __init__(self, requests_per_second: float = 2.0, burst: int = 10):
        self.rps = requests_per_second
        self.burst = burst
        self.tokens: Dict[str, float] = defaultdict(lambda: burst)
        self.last_update: Dict[str, float] = defaultdict(time.time)
        self._lock = threading.Lock()
    
    def _get_domain(self, url: str) -> str:
        try: return urlparse(url).netloc or "global"
        except: return "global"
    
    async def acquire(self, url: str = "global", timeout: float = 30) -> bool:
        domain = self._get_domain(url)
        start = time.time()
        
        while time.time() - start < timeout:
            with self._lock:
                now = time.time()
                elapsed = now - self.last_update[domain]
                self.tokens[domain] = min(self.burst, self.tokens[domain] + elapsed * self.rps)
                self.last_update[domain] = now
                
                if self.tokens[domain] >= 1:
                    self.tokens[domain] -= 1
                    return True
            
            await asyncio.sleep(0.1)
        
        return False
    
    def get_stats(self) -> Dict:
        return {domain: {'tokens': tokens, 'last': self.last_update[domain]}
                for domain, tokens in self.tokens.items()}

rate_limiter = RateLimiter()


# ═══════════════════════════════════════════════════════════════════════════════
# VISUALIZADOR v5.0 - Dashboard en tiempo real
# ═══════════════════════════════════════════════════════════════════════════════

class Visualizer:
    """Visualizador avanzado con métricas en tiempo real."""
    
    def __init__(self):
        self.console = console
        self.stats = {
            'tools_used': 0, 'files_created': 0, 'searches': 0,
            'pages_visited': 0, 'screenshots': 0, 'cache_hits': 0,
            'errors': 0, 'start_time': None
        }
        self.logs: deque = deque(maxlen=50)
        self.current_phase = ""
        self.current_step = ""
    
    def start(self):
        self.stats['start_time'] = datetime.now()
        self.stats = {k: 0 if isinstance(v, int) else v for k, v in self.stats.items()}
        self.stats['start_time'] = datetime.now()
    
    def log(self, msg: str, level: str = "info"):
        icons = {'info': 'ℹ️', 'success': '✅', 'warning': '⚠️', 'error': '❌', 
                 'tool': '🔧', 'browser': '🌐', 'search': '🔍', 'file': '📄',
                 'cache': '💾', 'step': '→', 'phase': '📋'}
        icon = icons.get(level, '•')
        timestamp = datetime.now().strftime("%H:%M:%S")
        log_entry = f"[dim]{timestamp}[/dim] {icon} {msg}"
        self.logs.append(log_entry)
        
        style = {'error': 'red', 'warning': 'yellow', 'success': 'green'}.get(level, 'white')
        self.console.print(f"    {icon} [{style}]{msg}[/{style}]")
    
    def phase(self, name: str, icon: str = "📋"):
        self.current_phase = name
        self.console.print(f"\n  [bold blue]{icon} {name}[/bold blue]")
    
    def step(self, desc: str):
        self.current_step = desc
        self.console.print(f"    [cyan]→[/cyan] {desc}")
    
    def tool(self, name: str, params: Dict):
        p_str = ", ".join(f"{k}={repr(v)[:20]}" for k, v in list(params.items())[:2])
        self.console.print(f"      [yellow]🔧 {name}[/yellow]({p_str})")
        self.stats['tools_used'] += 1
    
    def result(self, success: bool, msg: str):
        if success:
            self.console.print(f"        [green]✓[/green] {msg[:80]}")
        else:
            self.console.print(f"        [red]✗[/red] {msg[:80]}")
            self.stats['errors'] += 1
    
    def browser_action(self, action: str, url: str = ""):
        url_short = url[:50] + "..." if len(url) > 50 else url
        self.console.print(f"      [cyan]🌐 {action}[/cyan] {url_short}")
        if action == "navegando":
            self.stats['pages_visited'] += 1
    
    def screenshot_taken(self, path: str):
        self.console.print(f"      [magenta]📸 Screenshot:[/magenta] {Path(path).name}")
        self.stats['screenshots'] += 1
    
    def file_created(self, path: str):
        self.console.print(f"      [green]📄 Creado:[/green] {path}")
        self.stats['files_created'] += 1
    
    def search_result(self, query: str, count: int, cached: bool = False):
        cache_icon = " 💾" if cached else ""
        self.console.print(f"      [cyan]🔍 '{query}'[/cyan] → {count} resultados{cache_icon}")
        self.stats['searches'] += 1
        if cached: self.stats['cache_hits'] += 1
    
    def progress_bar(self, current: int, total: int, desc: str = ""):
        pct = (current / total) * 100 if total > 0 else 0
        filled = int(pct / 5)
        bar = "█" * filled + "░" * (20 - filled)
        self.console.print(f"      [{bar}] {pct:.0f}% {desc}")
    
    def show_summary(self, plan: TaskPlan):
        elapsed = (datetime.now() - self.stats['start_time']).total_seconds() if self.stats['start_time'] else 0
        
        self.console.print()
        self.console.print(Rule("[bold green]✅ PROCESO COMPLETADO[/bold green]"))
        
        # Tabla de estadísticas
        stats_table = Table(show_header=False, box=box.SIMPLE, padding=(0, 2))
        stats_table.add_column("Métrica", style="cyan")
        stats_table.add_column("Valor", style="green", justify="right")
        
        stats_table.add_row("⏱️ Duración", f"{elapsed:.1f}s")
        stats_table.add_row("📊 Progreso", f"{plan.progress:.0f}%")
        stats_table.add_row("🔧 Herramientas", str(self.stats['tools_used']))
        stats_table.add_row("📄 Archivos", str(self.stats['files_created']))
        stats_table.add_row("🔍 Búsquedas", str(self.stats['searches']))
        stats_table.add_row("🌐 Páginas", str(self.stats['pages_visited']))
        stats_table.add_row("📸 Screenshots", str(self.stats['screenshots']))
        stats_table.add_row("💾 Cache hits", str(self.stats['cache_hits']))
        
        if self.stats['errors'] > 0:
            stats_table.add_row("❌ Errores", str(self.stats['errors']))
        
        self.console.print(stats_table)
        
        # Árbol de fases
        tree = Tree("[bold]📋 Resumen de ejecución[/bold]")
        for phase in plan.phases:
            icon = "✅" if phase.is_complete else ("🔄" if phase.status == PhaseStatus.IN_PROGRESS else "⏳")
            branch = tree.add(f"{icon} {phase.icon} [cyan]{phase.name}[/cyan]")
            for step in phase.steps:
                s_icon = "✅" if step.status == PhaseStatus.COMPLETED else "❌"
                duration = f" ({step.duration:.1f}s)" if step.duration > 0 else ""
                branch.add(f"{s_icon} {step.description}{duration}")
        
        self.console.print(tree)

viz = Visualizer()

# ═══════════════════════════════════════════════════════════════════════════════
# SECURITY GUARD v5.0
# ═══════════════════════════════════════════════════════════════════════════════

class SecurityGuard:
    """Sistema de seguridad multicapa v5.0."""
    
    CRITICAL_PATTERNS = [
        r"rm\s+(-[rfv]+\s+)*/?$", r"rm\s+(-[rfv]+\s+)*/\*",
        r"rm\s+(-[rfv]+\s+)*/(home|etc|var|usr|bin|boot)",
        r":\(\)\s*\{\s*:\|:\s*&\s*\}\s*;", r"mkfs\.\w+",
        r"dd\s+if=/dev/(zero|random|urandom)\s+of=/dev/sd",
        r"chmod\s+(-R\s+)?(777|000)\s+/",
        r"curl.*\|\s*(bash|sh|python)", r"wget.*\|\s*(bash|sh|python)",
        r">\s*/dev/sd[a-z]", r"echo\s+.*>\s*/etc/(passwd|shadow)",
    ]
    
    # Shell injection patterns - command chaining and dangerous redirects
    INJECTION_PATTERNS = [
        r";\s*(rm|curl|wget|nc|netcat|bash|sh|python|perl|ruby)\s+",
        r"\|\s*(bash|sh|python|perl|ruby|nc|netcat)\s*",
        r"&&\s*(rm|curl|wget|nc|netcat|bash|sh|python)\s+",
        r"\|\|\s*(rm|curl|wget|nc|netcat|bash|sh)\s+",
        r"`.*rm\s+",
        r"\$\(.*rm\s+",
    ]
    
    SAFE_COMMANDS = {
        'ls', 'pwd', 'cd', 'cat', 'head', 'tail', 'less', 'more', 'echo',
        'date', 'whoami', 'hostname', 'uname', 'uptime', 'free', 'df', 'du',
        'ps', 'top', 'grep', 'find', 'wc', 'sort', 'uniq', 'cut', 'tr',
        'mkdir', 'touch', 'cp', 'mv', 'file', 'stat', 'tar', 'gzip', 'zip',
        'python', 'python3', 'pip', 'pip3', 'node', 'npm', 'git', 'curl', 'wget',
        'clear', 'history', 'which', 'whereis', 'type', 'man', 'help',
    }
    
    BLOCKED_DOMAINS = {
        'localhost', '127.0.0.1', '0.0.0.0', '::1',
        '192.168.', '10.', '172.16.', '172.17.', '172.18.', '172.19.',
    }
    
    def __init__(self, sandbox_root: str = None):
        self.sandbox_root = Path(sandbox_root or WORKSPACE).resolve()
        self.sandbox_root.mkdir(parents=True, exist_ok=True)
        self._critical = [re.compile(p, re.I) for p in self.CRITICAL_PATTERNS]
        self._injection = [re.compile(p, re.I) for p in self.INJECTION_PATTERNS]
        self.blocked_count = 0
        self.allowed_count = 0
    
    def analyze_command(self, cmd: str) -> SecurityAnalysis:
        cmd = cmd.strip()
        if not cmd:
            return SecurityAnalysis(cmd, True, ThreatLevel.SAFE, SecurityAction.ALLOW)
        
        # Críticos
        for p in self._critical:
            if p.search(cmd):
                self.blocked_count += 1
                return SecurityAnalysis(cmd, False, ThreatLevel.CRITICAL, 
                    SecurityAction.LOG_AND_BLOCK, ["🚨 COMANDO PELIGROSO BLOQUEADO"], 100.0)
        
        # Shell injection detection
        for p in self._injection:
            if p.search(cmd):
                self.blocked_count += 1
                return SecurityAnalysis(cmd, False, ThreatLevel.CRITICAL,
                    SecurityAction.LOG_AND_BLOCK, ["🚨 POSIBLE INYECCIÓN DE COMANDOS BLOQUEADA"], 100.0)
        
        # Comando base
        base = cmd.split()[0].split('/')[-1] if cmd.split() else ""
        if base in self.SAFE_COMMANDS:
            self.allowed_count += 1
            return SecurityAnalysis(cmd, True, ThreatLevel.SAFE, SecurityAction.ALLOW, [], 0.0)
        
        self.allowed_count += 1
        return SecurityAnalysis(cmd, True, ThreatLevel.LOW, SecurityAction.ALLOW, [], 10.0)
    
    def validate_url(self, url: str) -> Tuple[bool, str]:
        """Valida si una URL es segura para navegar."""
        try:
            parsed = urlparse(url)
            if not parsed.scheme in ['http', 'https']:
                return False, f"Esquema no permitido: {parsed.scheme}"
            
            host = parsed.netloc.lower()
            for blocked in self.BLOCKED_DOMAINS:
                if host.startswith(blocked) or host == blocked:
                    return False, f"Dominio bloqueado: {host}"
            
            return True, ""
        except Exception as e:
            return False, str(e)
    
    def validate_path(self, path: str) -> Dict[str, Any]:
        try:
            resolved = (self.sandbox_root / path).resolve() if not os.path.isabs(path) else Path(path).resolve()
            try:
                resolved.relative_to(self.sandbox_root)
                return {'allowed': True, 'path': str(resolved)}
            except ValueError:
                return {'allowed': False, 'reason': 'Fuera del sandbox'}
        except Exception as e:
            return {'allowed': False, 'reason': str(e)}


# ═══════════════════════════════════════════════════════════════════════════════
# NAVEGADOR WEB v5.0 - Computadora Virtual
# ═══════════════════════════════════════════════════════════════════════════════

class BrowserPool:
    """Pool de navegadores reutilizables para máximo rendimiento."""
    
    def __init__(self, max_browsers: int = MAX_CONCURRENT_BROWSERS):
        self.max_browsers = max_browsers
        self._playwright = None
        self._browser = None
        self._contexts: deque = deque(maxlen=max_browsers)
        self._available: asyncio.Queue = None
        self._lock = asyncio.Lock()
        self._initialized = False
        self.stats = {'created': 0, 'reused': 0, 'errors': 0}
    
    async def initialize(self):
        if self._initialized:
            return
        
        async with self._lock:
            if self._initialized:
                return
            
            if PLAYWRIGHT_OK:
                try:
                    from playwright.async_api import async_playwright
                    self._playwright = await async_playwright().start()
                    self._browser = await self._playwright.chromium.launch(
                        headless=True,
                        args=[
                            '--no-sandbox',
                            '--disable-setuid-sandbox',
                            '--disable-dev-shm-usage',
                            '--disable-gpu',
                            '--disable-web-security',
                            '--disable-features=IsolateOrigins,site-per-process'
                        ]
                    )
                    self._available = asyncio.Queue()
                    self._initialized = True
                    logger.info("BrowserPool inicializado con Playwright")
                except Exception as e:
                    logger.warning(f"No se pudo inicializar Playwright: {e}")
    
    async def get_context(self) -> Optional[Any]:
        """Obtiene un contexto de navegador del pool."""
        await self.initialize()
        
        if not self._browser:
            return None
        
        try:
            # Intentar reutilizar
            if not self._available.empty():
                ctx = await self._available.get()
                self.stats['reused'] += 1
                return ctx
            
            # Crear nuevo
            ctx = await self._browser.new_context(
                viewport={'width': 1920, 'height': 1080},
                user_agent=get_random_ua(),
                ignore_https_errors=True,
                java_script_enabled=True,
            )
            self.stats['created'] += 1
            return ctx
        except Exception as e:
            self.stats['errors'] += 1
            logger.error(f"Error obteniendo contexto: {e}")
            return None
    
    async def release_context(self, ctx):
        """Devuelve un contexto al pool."""
        if ctx and self._available:
            try:
                # Limpiar páginas
                for page in ctx.pages:
                    await page.close()
                await self._available.put(ctx)
            except:
                pass
    
    async def close(self):
        """Cierra todos los navegadores."""
        if self._browser:
            try:
                await self._browser.close()
            except: pass
        if self._playwright:
            try:
                await self._playwright.stop()
            except: pass
        self._initialized = False


class VirtualBrowser:
    """Navegador virtual completo - Computadora Virtual v5.0."""
    
    def __init__(self, security: SecurityGuard = None):
        self.security = security or SecurityGuard()
        self.pool = BrowserPool()
        self.engine = self._detect_engine()
        self.session_cookies: Dict[str, Dict] = {}
        self.history: List[str] = []
    
    def _detect_engine(self) -> BrowserEngine:
        if PLAYWRIGHT_OK:
            return BrowserEngine.PLAYWRIGHT
        elif SELENIUM_OK:
            return BrowserEngine.SELENIUM
        return BrowserEngine.HTTPX
    
    async def navigate(
        self,
        url: str,
        wait_for: str = "load",  # load, domcontentloaded, networkidle
        timeout: int = 30,
        screenshot: bool = False,
        extract_text: bool = True,
        wait_selector: str = None,
        scroll: bool = False,
        click_selector: str = None,
        fill_form: Dict[str, str] = None,
    ) -> WebPage:
        """
        Navega a una URL como un humano.
        
        Args:
            url: URL a visitar
            wait_for: Evento a esperar (load, domcontentloaded, networkidle)
            timeout: Timeout en segundos
            screenshot: Capturar screenshot
            extract_text: Extraer texto de la página
            wait_selector: Selector CSS a esperar
            scroll: Hacer scroll en la página
            click_selector: Selector a clickear
            fill_form: Formulario a llenar {selector: valor}
        """
        # Validar URL
        is_safe, reason = self.security.validate_url(url)
        if not is_safe:
            return WebPage(url=url, error=f"URL bloqueada: {reason}")
        
        # Verificar caché
        cache_key = f"browser:{url}:{extract_text}"
        cached = cache.get(cache_key)
        if cached and not screenshot:
            cached['cached'] = True
            return WebPage(**cached)
        
        # Rate limiting
        if not await rate_limiter.acquire(url):
            return WebPage(url=url, error="Rate limit excedido")
        
        start_time = time.time()
        viz.browser_action("navegando", url)
        
        # Usar el motor apropiado
        if self.engine == BrowserEngine.PLAYWRIGHT:
            result = await self._navigate_playwright(url, wait_for, timeout, screenshot, 
                                                     extract_text, wait_selector, scroll,
                                                     click_selector, fill_form)
        elif self.engine == BrowserEngine.SELENIUM:
            result = await self._navigate_selenium(url, timeout, screenshot, extract_text)
        else:
            result = await self._navigate_httpx(url, timeout, extract_text)
        
        result.load_time = time.time() - start_time
        
        # Guardar en caché si exitoso
        if not result.error and extract_text:
            cache_data = {
                'url': result.url, 'final_url': result.final_url,
                'title': result.title, 'text': result.text[:10000],
                'status_code': result.status_code, 'js_rendered': result.js_rendered
            }
            cache.set(cache_key, cache_data, CACHE_TTL)
        
        self.history.append(url)
        return result
    
    async def _navigate_playwright(
        self, url: str, wait_for: str, timeout: int, screenshot: bool,
        extract_text: bool, wait_selector: str, scroll: bool,
        click_selector: str, fill_form: Dict
    ) -> WebPage:
        """Navegación con Playwright."""
        ctx = await self.pool.get_context()
        if not ctx:
            return await self._navigate_httpx(url, timeout, extract_text)
        
        page = None
        try:
            page = await ctx.new_page()
            
            # Navegar
            response = await page.goto(url, wait_until=wait_for, timeout=timeout * 1000)
            
            # Esperar selector específico
            if wait_selector:
                try:
                    await page.wait_for_selector(wait_selector, timeout=10000)
                except: pass
            
            # Scroll
            if scroll:
                await page.evaluate("window.scrollTo(0, document.body.scrollHeight)")
                await asyncio.sleep(1)
                await page.evaluate("window.scrollTo(0, 0)")
            
            # Click
            if click_selector:
                try:
                    await page.click(click_selector)
                    await asyncio.sleep(1)
                except: pass
            
            # Fill form
            if fill_form:
                for selector, value in fill_form.items():
                    try:
                        await page.fill(selector, value)
                    except: pass
            
            # Obtener contenido
            final_url = page.url
            title = await page.title()
            html = await page.content()
            
            # Extraer texto
            text = ""
            if extract_text:
                text = await page.evaluate("""
                    () => {
                        const scripts = document.querySelectorAll('script, style, noscript');
                        scripts.forEach(s => s.remove());
                        return document.body.innerText || document.body.textContent || '';
                    }
                """)
            
            # Screenshot
            screenshot_path = None
            if screenshot:
                screenshot_path = str(SCREENSHOTS_DIR / f"screenshot_{uuid.uuid4().hex[:8]}.png")
                await page.screenshot(path=screenshot_path, full_page=True)
                viz.screenshot_taken(screenshot_path)
            
            # Links e imágenes
            links = await page.evaluate("Array.from(document.querySelectorAll('a[href]')).map(a => a.href).slice(0, 50)")
            images = await page.evaluate("Array.from(document.querySelectorAll('img[src]')).map(i => i.src).slice(0, 20)")
            
            # Metadata
            metadata = await page.evaluate("""
                () => {
                    const meta = {};
                    document.querySelectorAll('meta').forEach(m => {
                        const name = m.getAttribute('name') || m.getAttribute('property');
                        if (name) meta[name] = m.getAttribute('content');
                    });
                    return meta;
                }
            """)
            
            return WebPage(
                url=url, final_url=final_url, title=title,
                text=text[:50000], html=html[:100000],
                links=links, images=images, metadata=metadata,
                screenshot=screenshot_path,
                status_code=response.status if response else 0,
                js_rendered=True
            )
        
        except Exception as e:
            return WebPage(url=url, error=str(e))
        
        finally:
            if page:
                try: await page.close()
                except: pass
            await self.pool.release_context(ctx)
    
    async def _navigate_selenium(self, url: str, timeout: int, screenshot: bool, extract_text: bool) -> WebPage:
        """Navegación con Selenium (fallback)."""
        try:
            options = ChromeOptions()
            options.add_argument('--headless')
            options.add_argument('--no-sandbox')
            options.add_argument('--disable-dev-shm-usage')
            options.add_argument(f'--user-agent={get_random_ua()}')
            
            # Usar webdriver_manager si está disponible
            try:
                service = Service(ChromeDriverManager().install())
                driver = webdriver.Chrome(service=service, options=options)
            except:
                driver = webdriver.Chrome(options=options)
            
            driver.set_page_load_timeout(timeout)
            driver.get(url)
            
            title = driver.title
            html = driver.page_source
            text = driver.find_element(By.TAG_NAME, 'body').text if extract_text else ""
            
            screenshot_path = None
            if screenshot:
                screenshot_path = str(SCREENSHOTS_DIR / f"screenshot_{uuid.uuid4().hex[:8]}.png")
                driver.save_screenshot(screenshot_path)
            
            driver.quit()
            
            return WebPage(
                url=url, final_url=driver.current_url, title=title,
                text=text[:50000], html=html[:100000],
                screenshot=screenshot_path, js_rendered=True
            )
        except Exception as e:
            return WebPage(url=url, error=str(e))
    
    async def _navigate_httpx(self, url: str, timeout: int, extract_text: bool) -> WebPage:
        """Navegación básica con HTTPX (sin JS)."""
        if not HTTPX_OK:
            return WebPage(url=url, error="httpx no disponible")
        
        try:
            async with httpx.AsyncClient(timeout=timeout, follow_redirects=True) as client:
                response = await client.get(url, headers={'User-Agent': get_random_ua()})
                html = response.text
                
                title = ""
                text = ""
                links = []
                images = []
                metadata = {}
                
                if BS4_OK and extract_text:
                    soup = BeautifulSoup(html, 'lxml' if LXML_OK else 'html.parser')
                    
                    title_tag = soup.find('title')
                    title = title_tag.text.strip() if title_tag else ""
                    
                    for tag in soup(['script', 'style', 'noscript', 'nav', 'footer', 'header']):
                        tag.decompose()
                    
                    text = soup.get_text('\n', strip=True)
                    
                    links = [a.get('href', '') for a in soup.find_all('a', href=True)][:50]
                    images = [img.get('src', '') for img in soup.find_all('img', src=True)][:20]
                    
                    for meta in soup.find_all('meta'):
                        name = meta.get('name') or meta.get('property')
                        if name:
                            metadata[name] = meta.get('content', '')
                
                return WebPage(
                    url=url, final_url=str(response.url), title=title,
                    text=text[:50000], html=html[:100000],
                    links=links, images=images, metadata=metadata,
                    status_code=response.status_code, js_rendered=False
                )
        except Exception as e:
            return WebPage(url=url, error=str(e))
    
    async def search_google(self, query: str, num_results: int = 10) -> SearchResult:
        """Búsqueda en Google (requiere navegador para evitar bloqueos)."""
        url = f"https://www.google.com/search?q={quote_plus(query)}&num={num_results}"
        
        page = await self.navigate(url, wait_for="domcontentloaded", extract_text=True)
        if page.error:
            return SearchResult(query=query, results=[], source="google", error=page.error)
        
        results = []
        if BS4_OK:
            soup = BeautifulSoup(page.html, 'lxml' if LXML_OK else 'html.parser')
            for g in soup.find_all('div', class_='g')[:num_results]:
                link = g.find('a')
                title_div = g.find('h3')
                snippet_div = g.find('div', class_='VwiC3b')
                
                if link and title_div:
                    results.append({
                        'title': title_div.text,
                        'url': link.get('href', ''),
                        'snippet': snippet_div.text if snippet_div else ''
                    })
        
        return SearchResult(query=query, results=results, total=len(results), source="google")
    
    async def take_screenshot(self, url: str, full_page: bool = True) -> Optional[str]:
        """Captura screenshot de una página."""
        page = await self.navigate(url, screenshot=True)
        return page.screenshot
    
    async def close(self):
        """Cierra el navegador."""
        await self.pool.close()

# Instancia global del navegador
browser = VirtualBrowser()


# ═══════════════════════════════════════════════════════════════════════════════
# COMMAND EXECUTOR Y FILE MANAGER v5.0
# ═══════════════════════════════════════════════════════════════════════════════

class CommandExecutor:
    ALLOWED_INTERPRETERS: Dict[str, str] = {
        'python3': sys.executable,
        'python': sys.executable,
        'node': 'node',
        'bash': '/bin/bash',
        'sh': '/bin/sh'
    }
    
    SHELL_METACHARACTERS = frozenset('|><&;$`(){}[]\\!*?#~')
    
    SAFE_COMMANDS = frozenset([
        'ls', 'pwd', 'cat', 'echo', 'date', 'whoami', 'head', 'tail',
        'wc', 'grep', 'find', 'mkdir', 'touch', 'cp', 'mv', 'rm',
        'python3', 'python', 'pip', 'node', 'npm', 'git'
    ])
    
    def __init__(self, security: SecurityGuard = None, timeout: int = 30):
        self.security = security or SecurityGuard()
        self.timeout = timeout
        self.workdir = self.security.sandbox_root
        self.history: deque = deque(maxlen=500)
    
    def _contains_shell_metacharacters(self, cmd: str) -> bool:
        """Check if command contains shell metacharacters (potential injection)."""
        return any(c in cmd for c in self.SHELL_METACHARACTERS)
    
    def _validate_command_safety(self, cmd: str) -> Tuple[bool, str]:
        """Validate command is safe to execute. Returns (is_safe, reason)."""
        if not cmd or not cmd.strip():
            return False, "Empty command"
        
        if self._contains_shell_metacharacters(cmd):
            return False, f"Command contains shell metacharacters which could enable injection attacks. Use safe subprocess calls instead."
        
        try:
            args = shlex.split(cmd)
            if not args:
                return False, "Could not parse command"
            base_cmd = os.path.basename(args[0])
            if base_cmd not in self.SAFE_COMMANDS and args[0] != sys.executable:
                return False, f"Command '{base_cmd}' not in allowlist"
        except ValueError as e:
            return False, f"Invalid command syntax: {e}"
        
        return True, "OK"
    
    async def execute(self, cmd: str, timeout: int = None, _trusted: bool = False) -> ExecutionResult:
        """Execute command securely using subprocess exec (no shell).
        
        Args:
            cmd: Command string to execute
            timeout: Timeout in seconds
            _trusted: Internal flag - if True, bypasses metacharacter check (use only for internal calls)
        """
        start = time.time()
        timeout = timeout or self.timeout
        
        analysis = self.security.analyze_command(cmd)
        if analysis.action == SecurityAction.LOG_AND_BLOCK:
            return ExecutionResult(cmd, ExecutionStatus.BLOCKED, error_message="Bloqueado por seguridad")
        
        if not _trusted:
            is_safe, reason = self._validate_command_safety(cmd)
            if not is_safe:
                return ExecutionResult(cmd, ExecutionStatus.BLOCKED, error_message=f"Command rejected: {reason}")
        
        try:
            args = shlex.split(cmd)
            proc = await asyncio.create_subprocess_exec(
                *args, stdout=asyncio.subprocess.PIPE, stderr=asyncio.subprocess.PIPE,
                cwd=str(self.workdir), env=os.environ.copy())
            
            stdout, stderr = await asyncio.wait_for(proc.communicate(), timeout)
            result = ExecutionResult(cmd, ExecutionStatus.COMPLETED, proc.returncode,
                stdout.decode('utf-8', errors='replace'),
                stderr.decode('utf-8', errors='replace'),
                time.time() - start)
        except asyncio.TimeoutError:
            return ExecutionResult(cmd, ExecutionStatus.TIMEOUT, error_message=f"Timeout: {timeout}s")
        except Exception as e:
            result = ExecutionResult(cmd, ExecutionStatus.FAILED, error_message=str(e))
        
        self.history.append(result)
        return result
    
    async def run_script(self, code: str, interpreter: str = "python3", timeout: int = None) -> ExecutionResult:
        """Run a script with validated interpreter using safe subprocess execution."""
        interpreter_key = interpreter.lower().strip()
        if interpreter_key not in self.ALLOWED_INTERPRETERS:
            return ExecutionResult(
                f"{interpreter} [script]", ExecutionStatus.BLOCKED,
                error_message=f"Interpreter not allowed: {interpreter}. Allowed: {list(self.ALLOWED_INTERPRETERS.keys())}"
            )
        
        interpreter_path = self.ALLOWED_INTERPRETERS[interpreter_key]
        ext = {'python3': '.py', 'python': '.py', 'bash': '.sh', 'node': '.js', 'sh': '.sh'}.get(interpreter_key, '.py')
        script_path = self.workdir / f"_script_{uuid.uuid4().hex[:6]}{ext}"
        
        try:
            async with aiofiles.open(script_path, 'w') as f:
                await f.write(code)
            os.chmod(script_path, 0o755)
            
            start = time.time()
            timeout = timeout or self.timeout
            
            proc = await asyncio.create_subprocess_exec(
                interpreter_path, str(script_path),
                stdout=asyncio.subprocess.PIPE, stderr=asyncio.subprocess.PIPE,
                cwd=str(self.workdir), env=os.environ.copy())
            
            stdout, stderr = await asyncio.wait_for(proc.communicate(), timeout)
            result = ExecutionResult(
                f"{interpreter} {script_path.name}", ExecutionStatus.COMPLETED, proc.returncode,
                stdout.decode('utf-8', errors='replace'),
                stderr.decode('utf-8', errors='replace'),
                time.time() - start)
            self.history.append(result)
            return result
        except asyncio.TimeoutError:
            return ExecutionResult(f"{interpreter} [script]", ExecutionStatus.TIMEOUT, error_message=f"Timeout: {timeout}s")
        except Exception as e:
            return ExecutionResult(f"{interpreter} [script]", ExecutionStatus.FAILED, error_message=str(e))
        finally:
            try: script_path.unlink()
            except: pass


class FileManager:
    def __init__(self, security: SecurityGuard = None):
        self.security = security or SecurityGuard()
        self.root = self.security.sandbox_root
    
    def _resolve(self, p): return (self.root / p).resolve() if not os.path.isabs(p) else Path(p).resolve()
    def _validate(self, p): return self.security.validate_path(str(self._resolve(p)))
    
    async def read(self, path: str, encoding: str = 'utf-8') -> FileOperationResult:
        v = self._validate(path)
        if not v.get('allowed'): return FileOperationResult(False, 'read', path, error=v.get('reason'))
        try:
            async with aiofiles.open(self._resolve(path), 'r', encoding=encoding) as f:
                content = await f.read()
            return FileOperationResult(True, 'read', path, data=content, bytes_processed=len(content))
        except FileNotFoundError: return FileOperationResult(False, 'read', path, error="No encontrado")
        except Exception as e: return FileOperationResult(False, 'read', path, error=str(e))
    
    async def read_binary(self, path: str) -> FileOperationResult:
        v = self._validate(path)
        if not v.get('allowed'): return FileOperationResult(False, 'read_binary', path, error=v.get('reason'))
        try:
            async with aiofiles.open(self._resolve(path), 'rb') as f:
                return FileOperationResult(True, 'read_binary', path, data=await f.read())
        except Exception as e: return FileOperationResult(False, 'read_binary', path, error=str(e))
    
    async def write(self, path: str, content: str, create_dirs: bool = True) -> FileOperationResult:
        v = self._validate(path)
        if not v.get('allowed'): return FileOperationResult(False, 'write', path, error=v.get('reason'))
        resolved = self._resolve(path)
        try:
            if create_dirs: resolved.parent.mkdir(parents=True, exist_ok=True)
            async with aiofiles.open(resolved, 'w') as f:
                await f.write(content)
            return FileOperationResult(True, 'write', str(resolved), message=f"{len(content)} bytes", bytes_processed=len(content))
        except Exception as e: return FileOperationResult(False, 'write', path, error=str(e))
    
    async def write_binary(self, path: str, content: bytes) -> FileOperationResult:
        v = self._validate(path)
        if not v.get('allowed'): return FileOperationResult(False, 'write_binary', path, error=v.get('reason'))
        resolved = self._resolve(path)
        try:
            resolved.parent.mkdir(parents=True, exist_ok=True)
            async with aiofiles.open(resolved, 'wb') as f:
                await f.write(content)
            return FileOperationResult(True, 'write_binary', str(resolved), bytes_processed=len(content))
        except Exception as e: return FileOperationResult(False, 'write_binary', path, error=str(e))
    
    async def delete(self, path: str, recursive: bool = False) -> FileOperationResult:
        v = self._validate(path)
        if not v.get('allowed'): return FileOperationResult(False, 'delete', path, error=v.get('reason'))
        resolved = self._resolve(path)
        try:
            if resolved.is_file(): resolved.unlink()
            elif resolved.is_dir():
                if recursive: shutil.rmtree(resolved)
                else: resolved.rmdir()
            return FileOperationResult(True, 'delete', str(resolved))
        except Exception as e: return FileOperationResult(False, 'delete', path, error=str(e))
    
    async def list_dir(self, path: str = ".", pattern: str = "*", recursive: bool = False) -> FileOperationResult:
        v = self._validate(path)
        if not v.get('allowed'): return FileOperationResult(False, 'list', path, error=v.get('reason'))
        resolved = self._resolve(path)
        try:
            glob_fn = resolved.rglob if recursive else resolved.glob
            items = []
            for p in glob_fn(pattern):
                if p.name.startswith('.'): continue
                try:
                    stat = p.stat()
                    items.append({
                        'name': p.name,
                        'path': str(p.relative_to(self.root)) if p.is_relative_to(self.root) else str(p),
                        'is_dir': p.is_dir(),
                        'size': stat.st_size if p.is_file() else 0,
                        'modified': datetime.fromtimestamp(stat.st_mtime).isoformat()
                    })
                except: pass
            items.sort(key=lambda x: (not x['is_dir'], x['name'].lower()))
            return FileOperationResult(True, 'list', str(resolved), data={'items': items, 'count': len(items)})
        except Exception as e: return FileOperationResult(False, 'list', path, error=str(e))
    
    async def mkdir(self, path: str) -> FileOperationResult:
        v = self._validate(path)
        if not v.get('allowed'): return FileOperationResult(False, 'mkdir', path, error=v.get('reason'))
        try:
            self._resolve(path).mkdir(parents=True, exist_ok=True)
            return FileOperationResult(True, 'mkdir', path)
        except Exception as e: return FileOperationResult(False, 'mkdir', path, error=str(e))
    
    async def copy(self, src: str, dst: str) -> FileOperationResult:
        sv, dv = self._validate(src), self._validate(dst)
        if not sv.get('allowed') or not dv.get('allowed'):
            return FileOperationResult(False, 'copy', src, error="Ruta no permitida")
        try:
            sp, dp = self._resolve(src), self._resolve(dst)
            dp.parent.mkdir(parents=True, exist_ok=True)
            if sp.is_file(): shutil.copy2(sp, dp)
            else: shutil.copytree(sp, dp)
            return FileOperationResult(True, 'copy', str(dp))
        except Exception as e: return FileOperationResult(False, 'copy', src, error=str(e))
    
    async def move(self, src: str, dst: str) -> FileOperationResult:
        sv, dv = self._validate(src), self._validate(dst)
        if not sv.get('allowed') or not dv.get('allowed'):
            return FileOperationResult(False, 'move', src, error="Ruta no permitida")
        try:
            dp = self._resolve(dst)
            dp.parent.mkdir(parents=True, exist_ok=True)
            shutil.move(str(self._resolve(src)), str(dp))
            return FileOperationResult(True, 'move', str(dp))
        except Exception as e: return FileOperationResult(False, 'move', src, error=str(e))
    
    async def exists(self, path: str) -> bool:
        v = self._validate(path)
        return v.get('allowed') and self._resolve(path).exists()
    
    async def search(self, pattern: str, path: str = ".", content_search: str = None) -> FileOperationResult:
        result = await self.list_dir(path, pattern, recursive=True)
        if not result.success or not content_search:
            return result
        
        filtered = []
        for item in result.data.get('items', []):
            if item['is_dir']: continue
            try:
                r = await self.read(item['path'])
                if r.success and content_search.lower() in r.data.lower():
                    filtered.append(item)
            except: pass
        
        return FileOperationResult(True, 'search', path, data={'items': filtered, 'count': len(filtered)})

# ═══════════════════════════════════════════════════════════════════════════════
# DOCUMENT CREATOR v5.0 - Diseños Premium
# ═══════════════════════════════════════════════════════════════════════════════

class DocumentCreator:
    """Creador de documentos profesionales v5.0."""
    
    THEMES = {
        'professional': {'primary': '2E86AB', 'secondary': '2C3E50', 'accent': 'F18F01', 'bg': 'FFFFFF'},
        'modern': {'primary': '3498DB', 'secondary': '2C3E50', 'accent': '2ECC71', 'bg': 'F8F9FA'},
        'elegant': {'primary': '1A1A2E', 'secondary': '16213E', 'accent': 'E94560', 'bg': 'FFFFFF'},
        'creative': {'primary': '9B59B6', 'secondary': 'F1C40F', 'accent': 'E74C3C', 'bg': 'FFFFFF'},
        'minimal': {'primary': '333333', 'secondary': '666666', 'accent': '007AFF', 'bg': 'FFFFFF'},
        'dark': {'primary': 'BB86FC', 'secondary': '03DAC6', 'accent': 'CF6679', 'bg': '121212'},
    }
    
    def __init__(self, output_dir: str = None):
        self.output_dir = Path(output_dir or OUTPUT_DIR)
        self.output_dir.mkdir(parents=True, exist_ok=True)
    
    def _get_theme(self, name: str) -> Dict:
        return self.THEMES.get(name, self.THEMES['professional'])
    
    def _hex_to_rgb(self, hex_color: str) -> Tuple[int, int, int]:
        hex_color = hex_color.lstrip('#')
        return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
    
    async def create_pptx(
        self,
        title: str,
        slides: List[Dict],
        theme: str = "professional",
        author: str = "Agente IA v5.0",
        filename: str = None
    ) -> ToolResult:
        """Crea presentación PowerPoint con diseño premium."""
        if not PPTX_OK:
            return ToolResult(False, "pptx", error="python-pptx no instalado")
        
        try:
            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)
            
            colors = self._get_theme(theme)
            primary = RGBColor(*self._hex_to_rgb(colors['primary']))
            secondary = RGBColor(*self._hex_to_rgb(colors['secondary']))
            accent = RGBColor(*self._hex_to_rgb(colors['accent']))
            
            # Slide de título con gradiente visual
            title_slide = prs.slides.add_slide(prs.slide_layouts[6])
            
            # Fondo
            bg = title_slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
            bg.fill.solid()
            bg.fill.fore_color.rgb = primary
            bg.line.fill.background()
            
            # Decoración
            deco = title_slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(5.5), prs.slide_width, Inches(2))
            deco.fill.solid()
            deco.fill.fore_color.rgb = secondary
            deco.line.fill.background()
            
            # Título
            tb = title_slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(12.333), Inches(2))
            tf = tb.text_frame
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(60)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER
            
            # Subtítulo
            sub = title_slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(0.8))
            sf = sub.text_frame
            sp = sf.paragraphs[0]
            sp.text = f"Creado con Agente IA v5.0 • {datetime.now().strftime('%B %Y')}"
            sp.font.size = Pt(24)
            sp.font.color.rgb = RGBColor(200, 200, 200)
            sp.alignment = PP_ALIGN.CENTER
            
            # Slides de contenido
            for i, slide_data in enumerate(slides):
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                
                # Barra superior
                bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.1))
                bar.fill.solid()
                bar.fill.fore_color.rgb = primary
                bar.line.fill.background()
                
                # Título del slide
                st = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(1))
                stf = st.text_frame
                stp = stf.paragraphs[0]
                stp.text = slide_data.get('title', f'Slide {i+2}')
                stp.font.size = Pt(40)
                stp.font.bold = True
                stp.font.color.rgb = secondary
                
                # Línea decorativa bajo título
                line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.2), Inches(2), Inches(0.05))
                line.fill.solid()
                line.fill.fore_color.rgb = accent
                line.line.fill.background()
                
                content_top = Inches(1.5)
                
                # Bullets
                if slide_data.get('bullets'):
                    for j, bullet in enumerate(slide_data['bullets'][:10]):
                        bb = slide.shapes.add_textbox(Inches(0.7), content_top + Inches(j * 0.55), Inches(11.5), Inches(0.5))
                        bf = bb.text_frame
                        bp = bf.paragraphs[0]
                        bp.text = f"▸ {bullet}"
                        bp.font.size = Pt(22)
                        bp.font.color.rgb = secondary
                
                # Contenido texto
                elif slide_data.get('content'):
                    cb = slide.shapes.add_textbox(Inches(0.5), content_top, Inches(12.333), Inches(5))
                    cf = cb.text_frame
                    cf.word_wrap = True
                    cp = cf.paragraphs[0]
                    cp.text = slide_data['content']
                    cp.font.size = Pt(20)
                    cp.font.color.rgb = secondary
                
                # Dos columnas
                elif slide_data.get('columns'):
                    cols = slide_data['columns']
                    for ci, col in enumerate(cols[:2]):
                        x = Inches(0.5) if ci == 0 else Inches(6.9)
                        
                        # Título columna
                        ct = slide.shapes.add_textbox(x, content_top, Inches(5.8), Inches(0.6))
                        ctf = ct.text_frame
                        ctp = ctf.paragraphs[0]
                        ctp.text = col.get('title', '')
                        ctp.font.size = Pt(24)
                        ctp.font.bold = True
                        ctp.font.color.rgb = primary
                        
                        # Contenido columna
                        for bi, b in enumerate(col.get('items', [])[:6]):
                            bb = slide.shapes.add_textbox(x, content_top + Inches(0.7 + bi * 0.5), Inches(5.8), Inches(0.45))
                            bf = bb.text_frame
                            bp = bf.paragraphs[0]
                            bp.text = f"• {b}"
                            bp.font.size = Pt(18)
                
                # Número de página
                pn = slide.shapes.add_textbox(Inches(12.5), Inches(7), Inches(0.5), Inches(0.3))
                pnf = pn.text_frame
                pnp = pnf.paragraphs[0]
                pnp.text = str(i + 2)
                pnp.font.size = Pt(12)
                pnp.font.color.rgb = RGBColor(150, 150, 150)
            
            # Guardar
            filename = filename or f"presentacion_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
            filepath = self.output_dir / filename
            prs.save(str(filepath))
            
            return ToolResult(True, "pptx", data={'path': str(filepath), 'slides': len(slides) + 1},
                             message=f"Creado: {filename}", files_created=[str(filepath)])
        
        except Exception as e:
            return ToolResult(False, "pptx", error=str(e))
    
    async def create_docx(
        self,
        title: str,
        content: Union[str, List[Dict]],
        author: str = "Agente IA v5.0",
        theme: str = "professional",
        filename: str = None
    ) -> ToolResult:
        """Crea documento Word profesional."""
        if not DOCX_OK:
            return ToolResult(False, "docx", error="python-docx no instalado")
        
        try:
            doc = DocxDocument()
            colors = self._get_theme(theme)
            
            # Propiedades
            doc.core_properties.author = author
            doc.core_properties.title = title
            
            # Título
            title_para = doc.add_heading(title, 0)
            title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            
            # Línea decorativa
            doc.add_paragraph("━" * 60).alignment = WD_ALIGN_PARAGRAPH.CENTER
            
            # Metadatos
            meta = doc.add_paragraph()
            meta.alignment = WD_ALIGN_PARAGRAPH.CENTER
            run = meta.add_run(f"Autor: {author} | Fecha: {datetime.now().strftime('%d/%m/%Y')} | Creado con Agente IA v5.0")
            run.font.size = DP(10)
            run.font.italic = True
            
            doc.add_paragraph()
            
            # Contenido
            if isinstance(content, str):
                for para in content.split('\n\n'):
                    if para.strip():
                        p = doc.add_paragraph(para.strip())
                        p.paragraph_format.space_after = DP(12)
            else:
                for section in content:
                    if section.get('title'):
                        doc.add_heading(section['title'], level=section.get('level', 1))
                    
                    if section.get('content'):
                        p = doc.add_paragraph(section['content'])
                        p.paragraph_format.space_after = DP(12)
                    
                    for bullet in section.get('bullets', []):
                        doc.add_paragraph(bullet, style='List Bullet')
                    
                    # Tabla
                    table_data = section.get('table')
                    if table_data:
                        headers = table_data.get('headers', [])
                        rows = table_data.get('rows', [])
                        if headers:
                            table = doc.add_table(rows=1, cols=len(headers))
                            table.style = 'Table Grid'
                            
                            hdr_cells = table.rows[0].cells
                            for i, h in enumerate(headers):
                                hdr_cells[i].text = str(h)
                                hdr_cells[i].paragraphs[0].runs[0].bold = True
                            
                            for row in rows:
                                row_cells = table.add_row().cells
                                for i, cell in enumerate(row):
                                    if i < len(row_cells):
                                        row_cells[i].text = str(cell)
                            
                            doc.add_paragraph()
            
            filename = filename or f"documento_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
            filepath = self.output_dir / filename
            doc.save(str(filepath))
            
            return ToolResult(True, "docx", data={'path': str(filepath)},
                             message=f"Creado: {filename}", files_created=[str(filepath)])
        
        except Exception as e:
            return ToolResult(False, "docx", error=str(e))
    
    async def create_xlsx(
        self,
        title: str,
        sheets: List[Dict],
        theme: str = "professional",
        filename: str = None
    ) -> ToolResult:
        """Crea Excel con formato avanzado."""
        if not XLSX_OK:
            return ToolResult(False, "xlsx", error="openpyxl no instalado")
        
        try:
            wb = Workbook()
            wb.remove(wb.active)
            
            colors = self._get_theme(theme)
            
            header_font = Font(bold=True, color="FFFFFF", size=11)
            header_fill = PatternFill("solid", fgColor=colors['primary'])
            header_align = Alignment(horizontal="center", vertical="center")
            
            data_align = Alignment(horizontal="left", vertical="center")
            number_align = Alignment(horizontal="right", vertical="center")
            
            border = Border(
                left=Side(style='thin', color='CCCCCC'),
                right=Side(style='thin', color='CCCCCC'),
                top=Side(style='thin', color='CCCCCC'),
                bottom=Side(style='thin', color='CCCCCC')
            )
            
            alt_fill = PatternFill("solid", fgColor="F8F9FA")
            
            for sheet_data in sheets:
                ws = wb.create_sheet(title=sheet_data.get('name', 'Hoja')[:31])
                headers = sheet_data.get('headers', [])
                rows = sheet_data.get('rows', [])
                
                # Headers
                for col, header in enumerate(headers, 1):
                    cell = ws.cell(row=1, column=col, value=header)
                    cell.font = header_font
                    cell.fill = header_fill
                    cell.alignment = header_align
                    cell.border = border
                    ws.column_dimensions[get_column_letter(col)].width = max(15, len(str(header)) + 5)
                
                # Datos
                for row_idx, row_data in enumerate(rows, 2):
                    for col_idx, value in enumerate(row_data, 1):
                        cell = ws.cell(row=row_idx, column=col_idx, value=value)
                        cell.border = border
                        cell.alignment = number_align if isinstance(value, (int, float)) else data_align
                        
                        # Filas alternadas
                        if row_idx % 2 == 0:
                            cell.fill = alt_fill
                
                # Fórmulas
                for cell_ref, formula in sheet_data.get('formulas', {}).items():
                    ws[cell_ref] = formula
                
                # Gráfico
                if sheet_data.get('chart') and len(rows) > 0:
                    chart_config = sheet_data['chart']
                    chart_type = chart_config.get('type', 'bar')
                    
                    if chart_type == 'bar':
                        chart = BarChart()
                    elif chart_type == 'line':
                        chart = LineChart()
                    elif chart_type == 'pie':
                        chart = PieChart()
                    else:
                        chart = BarChart()
                    
                    chart.title = chart_config.get('title', 'Gráfico')
                    chart.style = 10
                    
                    data = Reference(ws, min_col=2, min_row=1, max_col=len(headers), max_row=len(rows)+1)
                    cats = Reference(ws, min_col=1, min_row=2, max_row=len(rows)+1)
                    chart.add_data(data, titles_from_data=True)
                    chart.set_categories(cats)
                    
                    ws.add_chart(chart, f"A{len(rows)+4}")
                
                ws.freeze_panes = 'A2'
            
            filename = filename or f"excel_{datetime.now().strftime('%Y%m%d_%H%M%S')}.xlsx"
            filepath = self.output_dir / filename
            wb.save(str(filepath))
            
            return ToolResult(True, "xlsx", data={'path': str(filepath), 'sheets': len(sheets)},
                             message=f"Creado: {filename}", files_created=[str(filepath)])
        
        except Exception as e:
            return ToolResult(False, "xlsx", error=str(e))


# ═══════════════════════════════════════════════════════════════════════════════
# WEB SEARCH ENGINE v5.0
# ═══════════════════════════════════════════════════════════════════════════════

class WebSearchEngine:
    """Motor de búsqueda web con múltiples fuentes y caché."""
    
    def __init__(self):
        self.browser = browser
    
    async def search(self, query: str, num_results: int = 10, use_browser: bool = False) -> SearchResult:
        cache_key = f"search:{query}:{num_results}"
        cached = cache.get(cache_key)
        if cached:
            viz.search_result(query, len(cached.get('results', [])), cached=True)
            return SearchResult(**cached, cached=True)
        
        result = await self._search_ddg_api(query, num_results)
        
        if len(result.results) < 3 and use_browser and PLAYWRIGHT_OK:
            result = await self._search_with_browser(query, num_results)
        
        if result.results:
            cache.set(cache_key, {'query': result.query, 'results': result.results,
                                  'total': result.total, 'source': result.source}, CACHE_TTL)
        
        viz.search_result(query, len(result.results))
        return result
    
    async def _search_ddg_api(self, query: str, num: int) -> SearchResult:
        if not HTTPX_OK:
            return SearchResult(query=query, source="ddg_api")
        try:
            async with httpx.AsyncClient(timeout=15) as client:
                r = await client.get("https://api.duckduckgo.com/",
                    params={'q': query, 'format': 'json', 'no_html': '1', 'skip_disambig': '1'})
                data = r.json()
                results = []
                if data.get('AbstractText'):
                    results.append({'title': data.get('Heading', query), 'snippet': data['AbstractText'][:400],
                                    'url': data.get('AbstractURL', ''), 'source': data.get('AbstractSource', 'DuckDuckGo')})
                for topic in data.get('RelatedTopics', [])[:num]:
                    if isinstance(topic, dict) and topic.get('Text'):
                        results.append({'title': topic.get('Text', '')[:100], 'snippet': topic.get('Text', ''),
                                       'url': topic.get('FirstURL', ''), 'source': 'DuckDuckGo'})
                return SearchResult(query=query, results=results[:num], total=len(results), source="ddg_api")
        except:
            return SearchResult(query=query, source="ddg_api")
    
    async def _search_with_browser(self, query: str, num: int) -> SearchResult:
        url = f"https://html.duckduckgo.com/html/?q={quote_plus(query)}"
        page = await self.browser.navigate(url, wait_for="domcontentloaded")
        if page.error:
            return SearchResult(query=query, source="ddg_browser")
        results = []
        if BS4_OK:
            soup = BeautifulSoup(page.html, 'lxml' if LXML_OK else 'html.parser')
            for result in soup.select('.result')[:num]:
                link = result.select_one('.result__a')
                snippet = result.select_one('.result__snippet')
                if link:
                    results.append({'title': link.text.strip(), 'url': link.get('href', ''),
                                   'snippet': snippet.text.strip() if snippet else '', 'source': 'DuckDuckGo'})
        return SearchResult(query=query, results=results, total=len(results), source="ddg_browser")
    
    async def search_news(self, query: str, num: int = 10) -> SearchResult:
        return await self.search(f"{query} news latest", num)

# ═══════════════════════════════════════════════════════════════════════════════
# HERRAMIENTAS v5.0
# ═══════════════════════════════════════════════════════════════════════════════

class BaseTool(ABC):
    name: str = "base"
    description: str = "Herramienta base"
    category: ToolCategory = ToolCategory.SYSTEM
    
    def __init__(self):
        self.enabled = True
        self.count = 0
        self.total_time = 0.0
    
    @abstractmethod
    async def execute(self, **kwargs) -> ToolResult: pass
    
    async def __call__(self, **kwargs) -> ToolResult:
        start = time.time()
        try:
            viz.tool(self.name, kwargs)
            result = await self.execute(**kwargs)
            result.execution_time = time.time() - start
            self.count += 1
            self.total_time += result.execution_time
            viz.result(result.success, result.message or str(result.data)[:60] if result.data else "OK")
            for f in result.files_created: viz.file_created(f)
            for s in result.screenshots: viz.screenshot_taken(s)
            return result
        except Exception as e:
            viz.result(False, str(e))
            return ToolResult(False, self.name, error=str(e))


class ShellTool(BaseTool):
    name = "shell"; description = "Ejecuta comandos shell de forma segura"
    category = ToolCategory.SYSTEM
    def __init__(self, ex=None): super().__init__(); self.ex = ex or CommandExecutor()
    async def execute(self, command: str, timeout: int = 30) -> ToolResult:
        r = await self.ex.execute(command, timeout)
        return ToolResult(r.success, self.name, data={'stdout': r.stdout, 'stderr': r.stderr, 'code': r.return_code},
                         message=r.stdout[:150] if r.stdout else "", error=r.error_message if not r.success else None)


class FileTool(BaseTool):
    name = "file"; description = "Operaciones de archivos: leer, escribir, listar, copiar, mover"
    category = ToolCategory.FILE
    def __init__(self, fm=None): super().__init__(); self.fm = fm or FileManager()
    async def execute(self, operation: str, path: str, content: str = None, destination: str = None, **kw) -> ToolResult:
        ops = {'read': lambda: self.fm.read(path), 'write': lambda: self.fm.write(path, content or ""),
               'delete': lambda: self.fm.delete(path, kw.get('recursive', False)),
               'list': lambda: self.fm.list_dir(path, kw.get('pattern', '*')),
               'mkdir': lambda: self.fm.mkdir(path), 'copy': lambda: self.fm.copy(path, destination or path + '_copy'),
               'move': lambda: self.fm.move(path, destination or path)}
        if operation not in ops: return ToolResult(False, self.name, error=f"Operación inválida: {operation}")
        r = await ops[operation]()
        if isinstance(r, bool): return ToolResult(True, self.name, data={'exists': r})
        fc = [path] if operation == 'write' and r.success else []
        return ToolResult(r.success, self.name, data=r.data, message=r.message, error=r.error, files_created=fc)


class PythonTool(BaseTool):
    name = "python"; description = "Ejecuta código Python"
    category = ToolCategory.DEVELOPMENT
    def __init__(self, ex=None): super().__init__(); self.ex = ex or CommandExecutor()
    async def execute(self, code: str, timeout: int = 60) -> ToolResult:
        imports = "import sys, os, json, math, random, datetime, re\n"
        r = await self.ex.run_script(imports + code, 'python3', timeout)
        return ToolResult(r.success, self.name, data={'output': r.stdout, 'errors': r.stderr},
                         message=r.stdout[:200] if r.stdout else "", error=r.error_message if not r.success else None)


class SearchTool(BaseTool):
    name = "search"; description = "Búsqueda web en tiempo real"
    category = ToolCategory.SEARCH
    def __init__(self): super().__init__(); self.engine = WebSearchEngine()
    async def execute(self, query: str, num_results: int = 10, use_browser: bool = False) -> ToolResult:
        r = await self.engine.search(query, num_results, use_browser)
        return ToolResult(len(r.results) > 0, self.name, data={'results': r.results, 'total': r.total, 'query': query},
                         message=f"{len(r.results)} resultados para '{query}'", cached=r.cached)


class BrowserTool(BaseTool):
    name = "browser"; description = "Navegador web completo - navega, screenshots, extrae contenido"
    category = ToolCategory.BROWSER
    async def execute(self, url: str, action: str = "navigate", screenshot: bool = False, 
                      wait_for: str = "load", scroll: bool = False, **kw) -> ToolResult:
        if action == "screenshot":
            path = await browser.take_screenshot(url)
            if path: return ToolResult(True, self.name, data={'screenshot': path}, message=f"Screenshot: {Path(path).name}", screenshots=[path])
            return ToolResult(False, self.name, error="No se pudo tomar screenshot")
        page = await browser.navigate(url, wait_for=wait_for, screenshot=screenshot, scroll=scroll)
        if page.error: return ToolResult(False, self.name, error=page.error)
        screenshots = [page.screenshot] if page.screenshot else []
        return ToolResult(True, self.name, data={'url': page.final_url or page.url, 'title': page.title, 
                         'text': page.text[:5000], 'status': page.status_code, 'js_rendered': page.js_rendered},
                         message=f"{page.title[:50]} ({page.status_code})", screenshots=screenshots)


class DocumentTool(BaseTool):
    name = "document"; description = "Crea documentos: PowerPoint, Word, Excel"
    category = ToolCategory.DOCUMENT
    def __init__(self, out=None): super().__init__(); self.creator = DocumentCreator(out)
    async def execute(self, doc_type: str, title: str, content: Any, theme: str = "professional", filename: str = None, **kw) -> ToolResult:
        doc_type = doc_type.lower()
        if doc_type == 'pptx':
            slides = content if isinstance(content, list) else [{'title': 'Contenido', 'content': str(content)}]
            return await self.creator.create_pptx(title, slides, theme, kw.get('author', 'Agente IA v5.0'), filename)
        elif doc_type == 'docx':
            return await self.creator.create_docx(title, content, kw.get('author', 'Agente IA v5.0'), theme, filename)
        elif doc_type == 'xlsx':
            sheets = content if isinstance(content, list) else [{'name': 'Datos', 'headers': [], 'rows': []}]
            return await self.creator.create_xlsx(title, sheets, theme, filename)
        return ToolResult(False, self.name, error=f"Tipo no soportado: {doc_type}")


class MessageTool(BaseTool):
    name = "message"; description = "Muestra mensajes al usuario"
    category = ToolCategory.COMMUNICATION
    async def execute(self, content: str, message_type: str = "info", title: str = None) -> ToolResult:
        icons = {'info': 'ℹ️', 'success': '✅', 'warning': '⚠️', 'error': '❌', 'tip': '💡'}
        icon = icons.get(message_type, '💬')
        if title: console.print(Panel(content, title=f"{icon} {title}", border_style="cyan"))
        else: console.print(f"    {icon} {content}")
        return ToolResult(True, self.name, message=content)


class ResearchTool(BaseTool):
    name = "research"; description = "Investigación profunda sobre un tema"
    category = ToolCategory.SEARCH
    def __init__(self): super().__init__(); self.search = WebSearchEngine()
    async def execute(self, topic: str, depth: str = "basic", max_sources: int = 5) -> ToolResult:
        viz.log(f"Investigando: {topic}", "search")
        sr = await self.search.search(topic, max_sources * 2)
        sources = []
        if depth in ["medium", "deep"] and sr.results:
            for i, result in enumerate(sr.results[:max_sources]):
                if result.get('url'):
                    page = await browser.navigate(result['url'])
                    if not page.error:
                        sources.append({'url': result['url'], 'title': page.title or result.get('title', ''),
                                       'content': page.text[:1000]})
                viz.progress_bar(i + 1, min(len(sr.results), max_sources), f"Fuente {i+1}")
        else:
            sources = [{'title': r.get('title', ''), 'snippet': r.get('snippet', ''), 'url': r.get('url', '')} 
                      for r in sr.results[:max_sources]]
        summary = f"Investigación sobre '{topic}': {len(sources)} fuentes analizadas."
        return ToolResult(True, self.name, data={'topic': topic, 'sources': sources, 'summary': summary},
                         message=f"Investigación: {len(sources)} fuentes")


class ToolRegistry:
    def __init__(self): self._tools: Dict[str, BaseTool] = {}
    def register(self, tool): self._tools[tool.name] = tool
    def get(self, name): return self._tools.get(name)
    async def execute(self, name: str, **params) -> ToolResult:
        tool = self.get(name)
        if not tool: return ToolResult(False, name, error=f"Herramienta no encontrada: {name}")
        return await tool(**params)
    def list_tools(self) -> List[str]: return list(self._tools.keys())
    def get_schemas(self) -> List[Dict]:
        return [{'name': t.name, 'description': t.description, 'category': t.category.value} for t in self._tools.values()]


# ═══════════════════════════════════════════════════════════════════════════════
# PLANIFICADOR v5.0
# ═══════════════════════════════════════════════════════════════════════════════

class TaskPlanner:
    PATTERNS = {
        'create_pptx': [r'crea.*(?:pptx?|powerpoint|presentaci)', r'make.*(?:ppt|presentation)'],
        'create_docx': [r'crea.*(?:docx?|word|documento)', r'make.*(?:doc|word)'],
        'create_xlsx': [r'crea.*(?:xlsx?|excel|hoja)', r'make.*(?:excel|spreadsheet)'],
        'research': [r'investig', r'research', r'analiz.*(?:tema|info)'],
        'browse': [r'naveg', r'visit', r'abre.*(?:url|p[aá]gina)', r'browse'],
        'screenshot': [r'screenshot', r'captura', r'pantallazo'],
        'search': [r'busc[ao]?\b', r'search\b', r'encuentra'],
        'file_create': [r'crea.*archivo', r'escrib.*archivo'],
        'file_read': [r'lee.*archivo', r'read.*file'],
        'file_list': [r'lista.*archivos', r'ls\b', r'dir\b'],
        'execute_code': [r'ejecut.*(?:c[oó]digo|python)', r'run.*code'],
        'install': [r'instal', r'pip\s+install'],
        'help': [r'ayuda', r'help', r'c[oó]mo'],
        'system': [r'sistema', r'status', r'info'],
    }
    
    def __init__(self):
        self._compiled = {k: [re.compile(p, re.I) for p in v] for k, v in self.PATTERNS.items()}
    
    def detect_intent(self, text: str) -> str:
        for intent, patterns in self._compiled.items():
            for p in patterns:
                if p.search(text): return intent
        return 'general'
    
    def extract_entities(self, text: str) -> Dict:
        return {'files': re.findall(r'[\w\-./]+\.\w{2,5}', text),
                'urls': re.findall(r'https?://[^\s<>"\']+', text),
                'quoted': re.findall(r'["\']([^"\']+)["\']', text),
                'topic': (re.search(r'(?:sobre|about|de)\s+["\']?([^"\',.!?]+)', text, re.I) or 
                         type('',(),{'group': lambda s,x: None})()).group(1)}
    
    async def create_plan(self, user_input: str, context=None) -> TaskPlan:
        intent = self.detect_intent(user_input)
        entities = self.extract_entities(user_input)
        viz.log(f"Intención: {intent}", "info")
        
        creators = {'create_pptx': self._pptx, 'create_docx': self._docx, 'create_xlsx': self._xlsx,
                   'research': self._research, 'browse': self._browse, 'screenshot': self._screenshot,
                   'search': self._search, 'file_create': self._file_create, 'file_read': self._file_read,
                   'file_list': self._file_list, 'execute_code': self._code, 'install': self._install,
                   'help': self._help, 'system': self._system}
        return await creators.get(intent, self._general)(user_input, entities)
    
    async def _pptx(self, text, e):
        topic = e.get('topic') or (e['quoted'][0] if e.get('quoted') else "Presentación")
        return TaskPlan(uuid.uuid4().hex[:8], f"Crear PowerPoint: {topic}", [
            Phase("research", "Investigación", "Recopilar info", "🔬", [
                Step("s1", f"Investigar {topic}", "research", {"topic": topic, "depth": "basic"})]),
            Phase("create", "Creación", "Crear presentación", "📊", [
                Step("s2", "Generar slides", "document", {"doc_type": "pptx", "title": topic,
                    "content": [{"title": "Introducción", "content": f"Sobre {topic}"},
                               {"title": "Puntos Principales", "bullets": ["Punto 1", "Punto 2", "Punto 3"]},
                               {"title": "Conclusión", "content": "Resumen final"}]})]),
            Phase("deliver", "Entrega", "Confirmar", "✅", [
                Step("s3", "Confirmar", "message", {"content": f"Presentación '{topic}' creada", "message_type": "success"})])])
    
    async def _docx(self, text, e):
        topic = e.get('topic') or "Documento"
        return TaskPlan(uuid.uuid4().hex[:8], f"Crear Word: {topic}", [
            Phase("create", "Creación", "Crear documento", "📝", [
                Step("s1", "Generar documento", "document", {"doc_type": "docx", "title": topic,
                    "content": [{"title": "Introducción", "level": 1, "content": f"Documento sobre {topic}."},
                               {"title": "Desarrollo", "level": 1, "content": "Contenido principal."}]})]),
            Phase("deliver", "Entrega", "Confirmar", "✅", [
                Step("s2", "Confirmar", "message", {"content": f"Documento '{topic}' creado", "message_type": "success"})])])
    
    async def _xlsx(self, text, e):
        topic = e.get('topic') or "Datos"
        return TaskPlan(uuid.uuid4().hex[:8], f"Crear Excel: {topic}", [
            Phase("create", "Creación", "Crear Excel", "📈", [
                Step("s1", "Generar Excel", "document", {"doc_type": "xlsx", "title": topic,
                    "content": [{"name": "Datos", "headers": ["A", "B", "C"],
                                "rows": [["D1", 100, 50], ["D2", 200, 75]],
                                "chart": {"type": "bar", "title": "Gráfico"}}]})])])
    
    async def _research(self, text, e):
        topic = e.get('topic') or text[:50]
        return TaskPlan(uuid.uuid4().hex[:8], f"Investigar: {topic}", [
            Phase("search", "Búsqueda", "Buscar info", "🔍", [Step("s1", f"Buscar {topic}", "search", {"query": topic})]),
            Phase("analyze", "Análisis", "Analizar", "🔬", [Step("s2", "Investigar", "research", {"topic": topic, "depth": "medium"})])])
    
    async def _browse(self, text, e):
        url = e['urls'][0] if e.get('urls') else "https://example.com"
        return TaskPlan(uuid.uuid4().hex[:8], f"Navegar: {url[:40]}", [
            Phase("browse", "Navegación", "Visitar", "🌐", [Step("s1", "Navegar", "browser", {"url": url, "screenshot": True})])])
    
    async def _screenshot(self, text, e):
        url = e['urls'][0] if e.get('urls') else "https://example.com"
        return TaskPlan(uuid.uuid4().hex[:8], f"Screenshot: {url[:40]}", [
            Phase("capture", "Captura", "Screenshot", "📸", [Step("s1", "Capturar", "browser", {"url": url, "action": "screenshot"})])])
    
    async def _search(self, text, e):
        query = text.replace('busca', '').strip()[:80]
        return TaskPlan(uuid.uuid4().hex[:8], f"Buscar: {query}", [
            Phase("search", "Búsqueda", "Buscar", "🔍", [Step("s1", f"Buscar: {query}", "search", {"query": query})])])
    
    async def _file_create(self, text, e):
        fn = e['files'][0] if e.get('files') else 'nuevo.txt'
        content = e['quoted'][0] if e.get('quoted') else "# Nuevo archivo\n"
        return TaskPlan(uuid.uuid4().hex[:8], f"Crear: {fn}", [
            Phase("create", "Crear", "Archivo", "📄", [Step("s1", f"Crear {fn}", "file", {"operation": "write", "path": fn, "content": content})])])
    
    async def _file_read(self, text, e):
        fn = e['files'][0] if e.get('files') else ''
        return TaskPlan(uuid.uuid4().hex[:8], f"Leer: {fn}", [
            Phase("read", "Leer", "Archivo", "📖", [Step("s1", f"Leer {fn}", "file", {"operation": "read", "path": fn})])])
    
    async def _file_list(self, text, e):
        return TaskPlan(uuid.uuid4().hex[:8], "Listar archivos", [
            Phase("list", "Listar", "Archivos", "📁", [Step("s1", "Listar", "file", {"operation": "list", "path": "."})])])
    
    async def _code(self, text, e):
        code = e.get('code', ['print("Hola v5!")'])[0] if 'code' in e else 'print("Hola v5!")'
        return TaskPlan(uuid.uuid4().hex[:8], "Ejecutar Python", [
            Phase("exec", "Ejecutar", "Python", "🐍", [Step("s1", "Ejecutar", "python", {"code": code})])])
    
    async def _install(self, text, e):
        m = re.search(r'install\s+(\S+)', text, re.I)
        pkg = m.group(1) if m else 'requests'
        return TaskPlan(uuid.uuid4().hex[:8], f"Instalar: {pkg}", [
            Phase("install", "Instalar", pkg, "📦", [Step("s1", f"pip install {pkg}", "shell", {"command": f"pip install {pkg} -q"})])])
    
    async def _help(self, text, e):
        help_msg = """🚀 **Agente IA v5.0** - Capacidades:
🌐 Navegador Web Virtual - Screenshots, extracción
📊 Documentos: PowerPoint, Word, Excel
🔍 Búsqueda e Investigación web
📁 Gestión de archivos
🐍 Ejecución de código Python

Ejemplos: "Crea presentación sobre IA", "Navega a url", "Investiga sobre X" """
        return TaskPlan(uuid.uuid4().hex[:8], "Ayuda", [
            Phase("help", "Ayuda", "Info", "❓", [Step("s1", "Mostrar", "message", {"content": help_msg, "message_type": "info", "title": "Ayuda"})])])
    
    async def _system(self, text, e):
        return TaskPlan(uuid.uuid4().hex[:8], "Sistema", [
            Phase("info", "Sistema", "Info", "💻", [Step("s1", "Info", "shell", {"command": "uname -a && python3 --version"})])])
    
    async def _general(self, text, e):
        return TaskPlan(uuid.uuid4().hex[:8], text[:60], [
            Phase("process", "Procesar", "Solicitud", "🔄", [Step("s1", "Procesar", "message", {"content": f"Procesando: {text[:60]}..."})])])

# ═══════════════════════════════════════════════════════════════════════════════
# AGENTE PRINCIPAL v5.0
# ═══════════════════════════════════════════════════════════════════════════════

@dataclass
class AgentConfig:
    name: str = "Agente IA v5.0"
    max_iterations: int = 100
    timeout: int = 60
    verbose: bool = True

class Agent:
    def __init__(self, config: AgentConfig = None):
        self.config = config or AgentConfig()
        self.state = AgentState.IDLE
        self.security = SecurityGuard()
        self.executor = CommandExecutor(self.security, self.config.timeout)
        self.files = FileManager(self.security)
        self.tools = ToolRegistry()
        self._register_tools()
        self.planner = TaskPlanner()
        self.current_plan: Optional[TaskPlan] = None
        self.history: List[Dict] = []
        self.iteration = 0
        viz.start()
    
    def _register_tools(self):
        self.tools.register(ShellTool(self.executor))
        self.tools.register(FileTool(self.files))
        self.tools.register(PythonTool(self.executor))
        self.tools.register(SearchTool())
        self.tools.register(BrowserTool())
        self.tools.register(DocumentTool(str(OUTPUT_DIR)))
        self.tools.register(MessageTool())
        self.tools.register(ResearchTool())
    
    async def run(self, user_input: str) -> str:
        self.iteration = 0
        self.state = AgentState.ANALYZING
        viz.start()
        
        try:
            if self.config.verbose:
                console.print()
                console.print(Rule(f"[bold blue]🚀 {self.config.name}[/bold blue]"))
                console.print(f"[dim]📝 {user_input[:80]}[/dim]\n")
            
            viz.phase("Planificación", "📋")
            self.current_plan = await self.planner.create_plan(user_input)
            if self.current_plan.phases:
                self.current_plan.phases[0].status = PhaseStatus.IN_PROGRESS
            
            if self.config.verbose:
                tree = Tree(f"[bold]📋 {self.current_plan.objective[:50]}[/bold]")
                for p in self.current_plan.phases:
                    branch = tree.add(f"{p.icon} [cyan]{p.name}[/cyan]")
                    for s in p.steps: branch.add(f"[dim]• {s.description}[/dim]")
                console.print(tree); console.print()
            
            while not self.current_plan.is_complete and self.iteration < self.config.max_iterations:
                self.iteration += 1
                phase = self.current_plan.get_current_phase()
                if not phase: break
                viz.phase(phase.name, phase.icon)
                step = phase.get_next_step()
                if not step: self.current_plan.advance(); continue
                
                self.state = AgentState.EXECUTING
                step.start()
                viz.step(step.description)
                result = await self.tools.execute(step.tool, **step.params)
                
                if result.success: step.complete(result.data)
                else: step.fail(result.error or "Error")
                
                self.history.append({'tool': step.tool, 'success': result.success})
                self.current_plan.advance()
            
            self.state = AgentState.DELIVERING
            response = self._build_response()
            if self.config.verbose: viz.show_summary(self.current_plan)
            return response
        except Exception as e:
            self.state = AgentState.ERROR
            return f"Error: {e}"
        finally:
            self.state = AgentState.IDLE
    
    def _build_response(self) -> str:
        if not self.current_plan: return "Sin plan."
        parts = [f"**{self.current_plan.objective}**"]
        files, results = [], []
        for p in self.current_plan.phases:
            for s in p.steps:
                if s.result and isinstance(s.result, dict):
                    if s.result.get('path'): files.append(s.result['path'])
                    if s.result.get('stdout'): results.append(s.result['stdout'][:150])
                    if s.result.get('output'): results.append(s.result['output'][:150])
                    if s.result.get('results'):
                        for r in s.result['results'][:3]: results.append(f"• {r.get('title', '')[:50]}")
        if files: parts.append("\n📁 **Archivos:** " + ", ".join(files))
        if results: parts.append("\n📊 **Resultados:**\n" + "\n".join(results[:5]))
        parts.append(f"\n✅ Completado en {self.iteration} pasos")
        return "\n".join(parts)
    
    async def execute_direct(self, tool_name: str, **params) -> ToolResult:
        return await self.tools.execute(tool_name, **params)
    
    def get_status(self) -> Dict:
        return {'name': self.config.name, 'version': VERSION, 'state': self.state.value,
                'tools': len(self.tools.list_tools()), 'browser': browser.engine.value}
    
    async def cleanup(self): await browser.close()

# ═══════════════════════════════════════════════════════════════════════════════
# INTERFAZ DE USUARIO
# ═══════════════════════════════════════════════════════════════════════════════

def print_banner():
    console.print("""
[bold blue]╔═════════════════════════════════════════════════════════════════════════════╗
║            🚀 AGENTE IA v5.0 - ENTERPRISE EDITION 🚀                        ║
╠═════════════════════════════════════════════════════════════════════════════╣
║  🌐 Navegador Virtual  |  📊 Documentos Pro  |  🔍 Búsqueda Web             ║
║  📸 Screenshots        |  💾 Caché Multinivel |  👁️ Visualización Real-Time║
╚═════════════════════════════════════════════════════════════════════════════╝[/bold blue]""")

async def demo():
    print_banner()
    console.print("\n[bold cyan]🎯 DEMO v5.0[/bold cyan]\n")
    agent = Agent(AgentConfig(verbose=True))
    for title, task in [("Sistema", "info del sistema"), ("Archivo", "crea archivo demo.txt con 'Hola v5!'"),
                        ("Búsqueda", "busca Python programming"), ("Listar", "lista archivos")]:
        console.print(f"\n[yellow]═══ {title} ═══[/yellow]")
        result = await agent.run(task)
        console.print(f"[dim]{result[:200]}[/dim]")
        await asyncio.sleep(0.5)
    console.print("\n[green]✅ Demo completada[/green]")
    await agent.cleanup()

async def interactive():
    print_banner()
    agent = Agent(AgentConfig(verbose=True))
    console.print(f"\n[green]✅ Agente v5.0 listo[/green] | Navegador: {browser.engine.value}")
    console.print("[dim]'help' para ayuda, 'exit' para salir[/dim]\n")
    
    try:
        while True:
            try:
                inp = Prompt.ask("[bold green]🚀 v5.0[/bold green]").strip()
                if not inp: continue
                cmd = inp.lower()
                if cmd in ['exit', 'salir', 'quit']: console.print("\n[yellow]👋 ¡Adiós![/yellow]"); break
                elif cmd in ['help', 'ayuda']: await agent.run("ayuda"); continue
                elif cmd == 'status':
                    for k, v in agent.get_status().items(): console.print(f"  {k}: {v}")
                    continue
                elif cmd == 'tools':
                    for t in agent.tools.get_schemas(): console.print(f"  • {t['name']}: {t['description'][:40]}")
                    continue
                elif cmd == 'clear': console.clear(); print_banner(); continue
                elif inp.startswith('!'):
                    r = await agent.execute_direct('shell', command=inp[1:])
                    console.print(r.data.get('stdout', '') if r.success else f"[red]{r.error}[/red]")
                    continue
                elif inp.startswith('py:'):
                    r = await agent.execute_direct('python', code=inp[3:])
                    console.print(r.data.get('output', '') if r.success else f"[red]{r.error}[/red]")
                    continue
                result = await agent.run(inp)
                console.print(f"\n[bold]📋 Resultado:[/bold]\n{result}\n")
            except KeyboardInterrupt:
                console.print("\n[yellow]Ctrl+C de nuevo para salir[/yellow]")
                try: await asyncio.sleep(1)
                except KeyboardInterrupt: break
            except Exception as e: console.print(f"[red]Error: {e}[/red]")
    finally:
        await agent.cleanup()

async def tests():
    console.print("\n[bold cyan]🧪 TESTS v5.0[/bold cyan]\n")
    passed = 0
    s = SecurityGuard()
    if s.analyze_command("ls").is_safe and not s.analyze_command("rm -rf /").is_safe: console.print("✅ Security"); passed += 1
    else: console.print("❌ Security")
    
    ex = CommandExecutor(s)
    r = await ex.execute("echo test")
    if r.success: console.print("✅ Executor"); passed += 1
    else: console.print("❌ Executor")
    
    fm = FileManager(s)
    await fm.write("test.txt", "v5")
    r = await fm.read("test.txt")
    if r.success and "v5" in r.data: console.print("✅ FileManager"); passed += 1
    else: console.print("❌ FileManager")
    await fm.delete("test.txt")
    
    agent = Agent(AgentConfig(verbose=False))
    r = await agent.run("ayuda")
    if r and len(r) > 20: console.print("✅ Agent"); passed += 1
    else: console.print("❌ Agent")
    await agent.cleanup()
    
    console.print(f"\n[bold]{passed}/4 tests OK[/bold]")

async def main():
    print_banner()
    console.print("\n[bold]Opciones:[/bold] 1=Interactivo 2=Demo 3=Tests 0=Salir")
    choice = Prompt.ask("[bold]Opción[/bold]", default="1")
    if choice == "1": await interactive()
    elif choice == "2": await demo()
    elif choice == "3": await tests()
    else: console.print("[yellow]👋 ¡Adiós![/yellow]")

if __name__ == "__main__":
    try: asyncio.run(main())
    except KeyboardInterrupt: console.print("\n[yellow]👋 ¡Adiós![/yellow]")
    except Exception as e: console.print(f"[red]Error: {e}[/red]"); traceback.print_exc()
